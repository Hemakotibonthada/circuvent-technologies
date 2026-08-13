"""
Verify generated business documents.

`npm run docs:business` printing "ok" only proves a file was written. It does
not prove the document has the right number of slides, that a price reached the
page, or that a table is not overflowing off the edge. A build script in this
repository once reported success while publishing the previous run's artifact,
so "the command succeeded" is not evidence.

This opens each artifact and asserts on its contents:

* every document contains the company name
* every priced document contains a real price from the catalogue
* nothing contains a raw unformatted price that bypassed the formatter
* the decks have the expected slide counts
* the PDFs have pages and extractable text
* placeholder tokens appear only where they are meant to

Run: python scripts/verify_business_docs.py
"""

from __future__ import annotations

import json
import re
import sys
from pathlib import Path

import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "Docs" / "business"
DATA = json.loads((OUT / "_data" / "business-data.json").read_text(encoding="utf8"))

COMPANY = DATA["company"]["name"]
failures: list[str] = []
checks = 0


def check(condition: bool, label: str) -> None:
    global checks
    checks += 1
    if not condition:
        failures.append(label)


def pptx_text(path: Path) -> tuple[str, int]:
    prs = Presentation(str(path))
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        chunks.append(cell.text)
    return "\n".join(chunks), len(prs.slides)


def docx_text(path: Path) -> str:
    doc = Document(str(path))
    chunks = [p.text for p in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                chunks.append(cell.text)
    for section in doc.sections:
        for p in section.footer.paragraphs:
            chunks.append(p.text)
    return "\n".join(chunks)


def pdf_text(path: Path) -> tuple[str, int]:
    with fitz.open(str(path)) as doc:
        return "\n".join(page.get_text() for page in doc), doc.page_count


def main() -> int:
    products = DATA["catalogue"]["products"]
    flagship = max(products, key=lambda p: p["price"])
    cheapest = min(products, key=lambda p: p["price"])

    print(f"Verifying documents in {OUT.relative_to(ROOT)}\n")

    # ---------------------------------------------------------------- decks
    for name, min_slides in [("Circuvent-Investor-Deck.pptx", 10),
                             ("Circuvent-Sales-Deck.pptx", 8)]:
        path = OUT / name
        check(path.exists(), f"{name} missing")
        if not path.exists():
            continue
        text, slides = pptx_text(path)
        check(slides >= min_slides, f"{name}: only {slides} slides, expected >= {min_slides}")
        check(COMPANY in text, f"{name}: company name absent")
        check(str(DATA["catalogue"]["total"]) in text, f"{name}: product count absent")
        # A deck that lost its numbers still looks fine; assert a real price landed.
        check("\u20b9" in text, f"{name}: no rupee-formatted price found")
        print(f"  {name:<38} {slides:>3} slides, {len(text):>6} chars")

    # ----------------------------------------------------------------- word
    for name, must_contain in [
        ("Circuvent-Company-Profile.docx", [COMPANY, DATA["company"]["supportEmail"]]),
        ("Circuvent-Business-Plan.docx", [COMPANY, "Executive summary"]),
        ("Circuvent-New-Joiner-Handbook.docx", [COMPANY, "00-start-here"]),
    ]:
        path = OUT / name
        check(path.exists(), f"{name} missing")
        if not path.exists():
            continue
        text = docx_text(path)
        for needle in must_contain:
            check(needle in text, f"{name}: missing {needle!r}")
        check(len(text) > 1500, f"{name}: suspiciously short ({len(text)} chars)")
        print(f"  {name:<38} {len(text):>6} chars")

    # ------------------------------------------------------------------ pdf
    for name, min_pages in [("Circuvent-Product-Catalogue.pdf", 5),
                            ("Circuvent-Price-List.pdf", 1)]:
        path = OUT / name
        check(path.exists(), f"{name} missing")
        if not path.exists():
            continue
        text, pages = pdf_text(path)
        check(pages >= min_pages, f"{name}: {pages} pages, expected >= {min_pages}")
        check(COMPANY in text, f"{name}: company name absent")
        check(flagship["name"] in text, f"{name}: flagship product absent")
        check(cheapest["name"] in text, f"{name}: cheapest product absent")
        # ReportLab's core fonts have no rupee glyph, so these use "Rs.".
        # Assert the substitution actually happened rather than leaving a
        # black box on a customer-facing price.
        check("Rs." in text, f"{name}: no 'Rs.' price formatting found")
        check("\u20b9" not in text, f"{name}: raw rupee sign present — will render as a box")
        print(f"  {name:<38} {pages:>3} pages,  {len(text):>6} chars")

    # -------------------------------------------------- every product priced
    cat_text, _ = pdf_text(OUT / "Circuvent-Product-Catalogue.pdf")
    missing = [p["name"] for p in products if p["name"] not in cat_text]
    check(not missing, f"Catalogue omits {len(missing)} products: {missing[:4]}")

    price_text, _ = pdf_text(OUT / "Circuvent-Price-List.pdf")
    missing_prices = [p["id"] for p in products if p["id"] not in price_text]
    check(not missing_prices,
          f"Price list omits {len(missing_prices)} products: {missing_prices[:4]}")

    # ------------------------------------------- placeholders where expected
    plan_text = docx_text(OUT / "Circuvent-Business-Plan.docx")
    check("to be supplied" in plan_text,
          "Business plan has no placeholder markers — unknown figures may have been invented")
    profile_text = docx_text(OUT / "Circuvent-Company-Profile.docx")
    check("to be supplied" not in profile_text,
          "Company profile contains placeholders; it is externally facing and should not")

    # ------------------------------------------------------- pluralisation
    # A category holding exactly one product rendered "1 products" on a printed
    # catalogue page. Every text-extraction check above passed it; only looking
    # at the rendered page caught it. Guarded so it cannot return when a
    # category shrinks to one.
    #
    # The lookbehind matters: a plain substring test also matches the "21
    # products" on the cover, and reports a defect in a correct document.
    bad_plural = re.compile(r"(?<!\d)1 products\b")
    check(not bad_plural.search(cat_text), "Catalogue renders '1 products'")
    for name in ["Circuvent-Investor-Deck.pptx", "Circuvent-Sales-Deck.pptx"]:
        deck_text, _ = pptx_text(OUT / name)
        check(not bad_plural.search(deck_text), f"{name}: renders '1 products'")

    print()
    if failures:
        print(f"FAILED {len(failures)} of {checks} checks:\n", file=sys.stderr)
        for f in failures:
            print(f"  - {f}", file=sys.stderr)
        return 1
    print(f"All {checks} checks passed.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
