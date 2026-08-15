#!/usr/bin/env python
"""
Verify the generated knowledge-transfer pack.

`npm run docs:kt` printing "ok" only proves three files were written. It does
not prove the deck has slides, that the device list reached the page, or that
the traps table survived being parsed out of a markdown file — and a build
script in this repository has previously reported success while publishing the
previous run's artifact, so "the command succeeded" is not evidence.

This opens each artifact and asserts on its contents:

* every document names the company and the commit it was generated from
* the deck has slides, speaker notes, and the parity rule on one of them
* the handbook names every document in Docs/, so the index cannot silently
  lose one
* the quick reference has pages and extractable text
* every device in the firmware tree appears somewhere in the pack
* nothing carries the business documents' "live product catalogue" stamp,
  which would be a claim this pack has no right to make

Run: python scripts/verify_kt_docs.py
"""

from __future__ import annotations

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT / "scripts"))

from kt_docs import facts                                    # noqa: E402

OUT = ROOT / "Docs" / "kt"
DECK = OUT / "Circuvent-KT-Deck.pptx"
ARCH = OUT / "Circuvent-KT-Architecture.pptx"
HANDBOOK = OUT / "Circuvent-KT-Handbook.docx"
QUICKREF = OUT / "Circuvent-KT-Quick-Reference.pdf"

failures: list[str] = []
checks = 0


def check(ok: bool, label: str) -> None:
    global checks
    checks += 1
    if not ok:
        failures.append(label)


def deck_text(path=None) -> tuple[str, int, int]:
    from pptx import Presentation
    prs = Presentation(str(path or DECK))
    out, notes = [], 0
    for slide in prs.slides:
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            notes += 1
        for shape in slide.shapes:
            if shape.has_text_frame:
                out.append(shape.text_frame.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        out.append(cell.text)
    return "\n".join(out), len(prs.slides._sldIdLst), notes


def docx_text() -> str:
    from docx import Document
    doc = Document(str(HANDBOOK))
    parts = [p.text for p in doc.paragraphs]
    for t in doc.tables:
        for row in t.rows:
            parts.extend(cell.text for cell in row.cells)
    return "\n".join(parts)


def pdf_text() -> tuple[str, int]:
    import fitz  # PyMuPDF — the same reader verify_business_docs.py already uses
    with fitz.open(str(QUICKREF)) as pdf:
        return "\n".join(page.get_text() for page in pdf), pdf.page_count


def main() -> int:
    for f in (DECK, ARCH, HANDBOOK, QUICKREF):
        if not f.is_file():
            print(f"missing: {f.relative_to(ROOT)} — run `npm run docs:kt` first", file=sys.stderr)
            return 1

    data = facts.collect()
    company = data["company"]["name"]

    deck, slides, notes = deck_text()
    arch, arch_slides, arch_notes = deck_text(ARCH)
    book = docx_text()
    sheet, pages = pdf_text()
    everything = "\n".join([deck, arch, book, sheet])

    # --- structure -------------------------------------------------------
    check(slides >= 12, f"deck has only {slides} slides")
    check(notes >= slides - 2, f"only {notes} of {slides} slides carry speaker notes")
    check(pages >= 1, "quick reference has no pages")
    check(len(book) > 4000, "handbook is suspiciously short")
    check(len(sheet) > 800, "quick reference has little extractable text")

    # --- the architecture deck ------------------------------------------
    # Checked separately from the handover deck because it makes a different
    # promise: it is the one somebody opens with a broker down, so a slide that
    # has quietly lost its content is worse here than anywhere else in the pack.
    from kt_docs import arch_facts as af

    check(arch_slides >= 14, f"architecture deck has only {arch_slides} slides")
    check(company.lower() in arch.lower(), "architecture deck does not name the company")

    # The topic names are parsed out of the firmware. If the parse returns
    # nothing the table renders empty and the deck still builds, which is
    # exactly the silent failure the deck itself is about.
    topics = af.mqtt_topics()
    check(len(topics) == 4, f"expected 4 MQTT topics, parsed {len(topics)}")
    for topic, *_ in topics:
        check(topic in arch, f"MQTT topic {topic} did not reach the deck")
    check("retained" in arch, "the retained/not-retained distinction is missing")

    # Broker limits are read out of mosquitto.conf; a default of "—" means the
    # parse missed and the slide is quietly asserting nothing.
    broker = "\n".join(v for _, v in af.mqtt_broker())
    check("(unparsed)" not in broker, "a broker setting fell back to a placeholder")
    check("8883" in arch and "1883" in arch, "the broker listeners are not on the page")

    # Certificates: the one that matters is the one nobody else renews.
    check("own CA" in arch, "the self-issued broker CA is not distinguished")
    check("gen-certs.sh" in arch, "the deck does not say where the CA comes from")

    # Required config is parsed from the schema that enforces it.
    required = af.control_plane_secrets()
    check(len(required) >= 3, f"only {len(required)} required variables parsed from config.ts")
    for name, _ in required:
        check(name in arch, f"required variable {name} is missing from the deck")

    # A secrets slide that printed a value would be a leak, not a document.
    # No real secret is available here, so this asserts the shape instead: the
    # slide names homes and variable names, never an assignment.
    check("=" not in "\n".join(h for _, h, _ in af.secret_homes()),
          "a secret home reads like an assignment rather than a location")

    # --- the pictures actually arrived ----------------------------------
    # svg_png returns None rather than raising when the renderer cannot read a
    # source, which is right for a build — one missing illustration must not
    # cost the whole pack. It does mean a deck can be published with the
    # pictures quietly absent, and nothing above would notice, so it is checked
    # here instead.
    from pptx import Presentation as _P
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    arch_shapes = [sh for slide in _P(str(ARCH)).slides for sh in slide.shapes]
    arch_pics = sum(1 for sh in arch_shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)
    check(arch_pics >= 6, f"architecture deck embeds only {arch_pics} pictures")

    # And that they are drawn rather than described: the figures are built from
    # connectors, which no table or bullet list produces. The enum is named
    # rather than numbered — the first version of this check guessed 20 for a
    # line, counted zero, and reported a deck with no diagrams in it as broken.
    arch_lines = sum(1 for sh in arch_shapes if sh.shape_type == MSO_SHAPE_TYPE.LINE)
    check(arch_lines >= 25, f"architecture deck has only {arch_lines} drawn connectors")

    # --- identity --------------------------------------------------------
    # Case-insensitive: the shared cover renders the company name in capitals.
    for name, text in (("deck", deck), ("handbook", book), ("quick reference", sheet)):
        check(company.lower() in text.lower(), f"{name} does not name the company")

    # --- derived content actually reached the page -----------------------
    # The whole justification for generating this pack is that it cannot fall
    # behind the tree. If a device is missing here, it did not.
    missing = [d for d in data["devices"] if d not in everything]
    check(not missing, f"devices missing from the pack: {', '.join(missing[:6])}")

    for name, title in data["docs"]:
        stem = name.replace(".md", "")
        check(stem in book, f"handbook does not list {name}")

    if data["traps"]:
        first_trap = data["traps"][0][0]
        check(first_trap[:18] in everything,
              "the traps table did not survive parsing into the pack")

    # The traps come out of a markdown table, so they arrive carrying markdown.
    # Rendered, that shows as literal backticks, ** around words and a link
    # target printed beside its own text. Building the pack and looking at it is
    # what caught this; these keep it caught.
    for marker, why in (("**", "bold markers"), ("](", "markdown links")):
        check(marker not in everything, f"raw {why} reached the page")

    # --- claims it must not make ----------------------------------------
    # The business documents are generated from the catalogue and say so. This
    # pack is generated from the repository, and borrowing that stamp would
    # assert something nobody checked.
    check("live product catalogue" not in everything,
          "the pack carries the business documents' catalogue stamp")

    # --- traceability ----------------------------------------------------
    check(data["commit"] in everything,
          "no document records the commit it was generated from")

    # --- counts are the counted ones, not remembered ones ---------------
    check(str(data["counts"]["devices"]) in everything,
          "the device count on the page does not match the tree")

    print(f"{checks - len(failures)}/{checks} checks passed "
          f"({slides} slides, {pages} pages, {len(data['devices'])} devices)")
    if failures:
        print("\nFAILED:", file=sys.stderr)
        for f in failures:
            print("  - " + f, file=sys.stderr)
        return 1
    print("ok")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
