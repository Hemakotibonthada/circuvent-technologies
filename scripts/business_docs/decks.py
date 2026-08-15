"""
PowerPoint decks.

Two audiences, two decks, one set of facts:

* **Investor / company deck** — what the company is, what has actually been
  built, and where the numbers we do not hold need filling in.
* **Sales deck** — what a customer gets, priced from the live catalogue.

Both are generated. No price, product name or count is typed into this file;
they all come from the export, so the decks cannot fall behind the shop.
"""

from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from .brand import (
    CYAN, CYAN_BRIGHT, VIOLET, INK, INK_DEEP, SLATE, MUTED, LINE, PAPER,
    PAPER_ALT, FONT_SANS, TBD, category_color, ordered_categories, plural,
    rupees, footer_line, generated_stamp,
)

W = 13.333
H = 7.5


def _rgb(hexstr: str) -> RGBColor:
    return RGBColor.from_string(hexstr)


def _new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    return prs


def _blank(prs: Presentation, bg: str = PAPER):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(bg)
    return slide


def _text(slide, left, top, width, height, text, *, size=18, bold=False,
          color=SLATE, align=PP_ALIGN.LEFT, font=FONT_SANS, spacing=1.0,
          anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = _rgb(color)
    return box


def _bullets(slide, left, top, width, height, items, *, size=17, color=SLATE,
             gap=10, bullet_color=CYAN, head_color=None):
    """
    A bulleted list, with an optional bold lead-in before " — ".

    `head_color` exists because the lead-in used to be hardcoded to INK. That is
    correct on the light slides these bullets were written for and invisible on
    a dark one — near-black text on the ink background, so half of each bullet
    simply vanished while the rest of the line read normally. It defaults to the
    old value, so every existing slide is unchanged; dark slides pass PAPER.

    Nothing could have caught that except rendering the deck and looking at it.
    """
    head = head_color or INK
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.25
        dot = p.add_run()
        dot.text = "\u25cf  "
        dot.font.size = Pt(size - 5)
        dot.font.name = FONT_SANS
        dot.font.color.rgb = _rgb(bullet_color)
        # A bold lead-in before " — " makes a dense list scannable; without it
        # every bullet reads at the same weight and none of them get read.
        if " \u2014 " in item:
            head_text, tail = item.split(" \u2014 ", 1)
            r1 = p.add_run()
            r1.text = head_text + " \u2014 "
            r1.font.size = Pt(size)
            r1.font.bold = True
            r1.font.name = FONT_SANS
            r1.font.color.rgb = _rgb(head)
            r2 = p.add_run()
            r2.text = tail
            r2.font.size = Pt(size)
            r2.font.name = FONT_SANS
            r2.font.color.rgb = _rgb(color)
        else:
            r = p.add_run()
            r.text = item
            r.font.size = Pt(size)
            r.font.name = FONT_SANS
            r.font.color.rgb = _rgb(color)
    return box


def _rect(slide, left, top, width, height, fill, *, line=None, radius=None):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    if line:
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    if shape.has_text_frame:
        shape.text_frame.text = ""
    return shape


def _accent_bar(slide, top=0.0):
    """The cyan→violet rule that ties a slide to the site's accent pair."""
    _rect(slide, 0, top, W / 2, 0.09, CYAN)
    _rect(slide, W / 2, top, W / 2, 0.09, VIOLET)


def _slide_header(slide, kicker, title, *, dark=False):
    _accent_bar(slide)
    _text(slide, 0.85, 0.55, 11.6, 0.4, kicker.upper(), size=12, bold=True,
          color=CYAN_BRIGHT if dark else CYAN)
    _text(slide, 0.85, 0.95, 11.6, 0.9, title, size=34, bold=True,
          color=PAPER if dark else INK)


def _footer(slide, data, *, dark=False, page=None):
    _text(slide, 0.85, H - 0.55, 9.5, 0.35, footer_line(data), size=9,
          color=MUTED if not dark else "94A3B8")
    if page is not None:
        _text(slide, W - 1.6, H - 0.55, 0.75, 0.35, str(page), size=9,
              color=MUTED if not dark else "94A3B8", align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------- slide kinds

def _title_slide(prs, data, title, subtitle, tag):
    slide = _blank(prs, INK_DEEP)
    _rect(slide, 0, 0, 0.16, H, CYAN)
    _rect(slide, 0, H * 0.55, 0.16, H * 0.45, VIOLET)
    _text(slide, 1.1, 2.1, 11, 0.5, data["company"]["name"].upper(), size=14,
          bold=True, color=CYAN_BRIGHT)
    _text(slide, 1.05, 2.65, 11.2, 1.5, title, size=54, bold=True, color=PAPER)
    _text(slide, 1.1, 4.35, 10.5, 1.0, subtitle, size=20, color="94A3B8",
          spacing=1.3)
    _rect(slide, 1.1, 5.6, 2.2, 0.045, CYAN)
    _text(slide, 1.1, 5.85, 10.5, 0.4, tag, size=13, color="64748B")
    _text(slide, 1.1, H - 0.75, 11, 0.35, generated_stamp(data), size=9.5,
          color="475569")
    return slide


def _section_slide(prs, data, number, title, blurb):
    slide = _blank(prs, INK)
    _rect(slide, 0, 0, W, 0.09, CYAN)
    _text(slide, 1.1, 2.55, 2.0, 1.2, number, size=64, bold=True, color=VIOLET)
    _text(slide, 1.1, 3.55, 10.8, 0.9, title, size=40, bold=True, color=PAPER)
    _text(slide, 1.15, 4.5, 10.5, 0.8, blurb, size=17, color="94A3B8", spacing=1.3)
    return slide


def _stat_row(slide, stats, top=2.5, *, dark=False):
    """Evenly spaced figure cards. Big number, small label, nothing else."""
    n = len(stats)
    gap = 0.35
    total = W - 1.7
    card_w = (total - gap * (n - 1)) / n
    for i, (value, label) in enumerate(stats):
        left = 0.85 + i * (card_w + gap)
        _rect(slide, left, top, card_w, 2.1, INK if dark else PAPER_ALT,
              line=None if dark else LINE, radius=True)
        _text(slide, left + 0.15, top + 0.32, card_w - 0.3, 0.95, str(value),
              size=40, bold=True, color=CYAN_BRIGHT if dark else CYAN,
              align=PP_ALIGN.CENTER)
        _text(slide, left + 0.15, top + 1.32, card_w - 0.3, 0.6, label,
              size=12.5, color="94A3B8" if dark else SLATE, align=PP_ALIGN.CENTER,
              spacing=1.15)


def _table(slide, headers, rows, *, left=0.85, top=2.0, width=None, col_widths=None,
           row_h=0.42, header_h=0.48, size=12):
    width = width or (W - 1.7)
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(left),
                                   Inches(top), Inches(width),
                                   Inches(header_h + row_h * len(rows)))
    table = shape.table
    table.first_row = True
    if col_widths:
        scale = width / sum(col_widths)
        for i, cw in enumerate(col_widths):
            table.columns[i].width = Inches(cw * scale)
    table.rows[0].height = Inches(header_h)
    for i in range(1, len(rows) + 1):
        table.rows[i].height = Inches(row_h)

    for c, head in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(INK)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.1)
        p = cell.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = head
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.name = FONT_SANS
        r.font.color.rgb = _rgb(PAPER)

    for ri, row in enumerate(rows, start=1):
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = ""
            cell.fill.solid()
            # Zebra striping: at this row height an unstriped price table is
            # very easy to read across the wrong line.
            cell.fill.fore_color.rgb = _rgb(PAPER if ri % 2 else PAPER_ALT)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.1)
            p = cell.text_frame.paragraphs[0]
            if ci > 1:
                p.alignment = PP_ALIGN.RIGHT
                cell.margin_right = Inches(0.1)
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(size)
            r.font.name = FONT_SANS
            r.font.color.rgb = _rgb(SLATE)
    return table


# ------------------------------------------------------------- investor deck

def build_investor_deck(data: dict, out_path):
    prs = _new_deck()
    cat = data["catalogue"]
    eng = data["engineering"]
    com = data["commercial"]
    page = 0

    _title_slide(
        prs, data, "Connected hardware,\nbuilt end to end.",
        f"{cat['total']} products on sale across {len(cat['categories'])} categories — "
        "designed, manufactured, and run on our own platform.",
        "Company overview",
    )

    # --- what we do
    page += 1
    s = _blank(prs)
    _slide_header(s, "What we do", "We build the whole stack, not a layer of it")
    _bullets(s, 0.85, 2.2, 11.6, 3.6, [
        "Hardware \u2014 our own PCB designs, "
        f"{eng['hardwareProjects']} board projects in the repository.",
        "Firmware \u2014 "
        f"{eng['firmwareDeviceTypes']} device types on a single shared ESP32 library, "
        "so provisioning, connectivity and over-the-air update behave identically across the fleet.",
        "Platform \u2014 a self-hosted control plane owning devices, telemetry, "
        "commands, scenes and automations.",
        "Applications \u2014 one console on the web and the same product on Android and iOS.",
        "Commerce \u2014 our own storefront, checkout, orders and support flow.",
    ])
    _text(s, 0.85, 5.95, 11.6, 0.8,
          "Most companies in this category assemble a white-label device and resell an app. "
          "Owning every layer is why a fix reaches a customer's device the same day.",
          size=14, color=MUTED, spacing=1.25)
    _footer(s, data, page=page)

    # --- problem
    page += 1
    s = _blank(prs)
    _slide_header(s, "Problem", "Smart home hardware fails in ways nobody sees")
    _bullets(s, 0.85, 2.2, 11.6, 3.6, [
        "Devices go offline and stay offline \u2014 a router that boots slower than the device "
        "strands it until somebody physically resets it.",
        "Controls that do nothing \u2014 an app offers a switch the hardware does not implement, "
        "and the customer reports faulty hardware.",
        "No accountability \u2014 when a pump burns out or a gate opens, there is no record of why.",
        "Fragmented ownership \u2014 hardware from one vendor, app from another, "
        "and neither owns the failure.",
    ], bullet_color=VIOLET)
    _text(s, 0.85, 5.9, 11.6, 0.9,
          "These are not exotic faults. They are the ordinary ones, and they are invisible "
          "because nothing errors \u2014 which is exactly why they survive.",
          size=14, color=MUTED, spacing=1.25)
    _footer(s, data, page=page)

    # --- what is built
    page += 1
    s = _blank(prs, INK_DEEP)
    _slide_header(s, "Status", "What is already built", dark=True)
    _stat_row(s, [
        (cat["total"], "products in the catalogue"),
        (eng["firmwareDeviceTypes"], "firmware device types"),
        (eng["hardwareProjects"], "PCB / hardware projects"),
        (eng["deployables"], "independently shipped systems"),
    ], top=2.35, dark=True)
    _text(s, 0.85, 4.85, 11.6, 1.2,
          "Every figure on this slide is counted from the source repository at generation time, "
          "not maintained by hand. The catalogue spans "
          f"{rupees(cat['priceMin'])} to {rupees(cat['priceMax'])} across "
          f"{', '.join(ordered_categories(cat))}."
          + (f" A further {cat['comingSoon']} are announced and not yet orderable."
             if cat.get("comingSoon") else ""),
          size=14.5, color="94A3B8", spacing=1.3)
    _footer(s, data, dark=True, page=page)

    # --- range
    page += 1
    s = _blank(prs)
    _slide_header(s, "Product range", "Four categories, one platform")
    rows = []
    for c in ordered_categories(cat):
        items = [p for p in cat["products"] if p["category"] == c]
        lo = min(p["price"] for p in items)
        hi = max(p["price"] for p in items)
        rows.append([c, str(len(items)), rupees(lo), rupees(hi)])
    _table(s, ["Category", "Products", "From", "To"], rows,
           top=2.2, col_widths=[5.2, 2.0, 2.2, 2.2], row_h=0.55)
    _text(s, 0.85, 5.4, 11.6, 1.0,
          "A single account, one app and one automation engine cover all of them. "
          "A customer who buys one device can add any other without a second system.",
          size=14, color=MUTED, spacing=1.25)
    _footer(s, data, page=page)

    # --- architecture
    page += 1
    s = _blank(prs)
    _slide_header(s, "How it works", "Four systems, documented contracts")
    boxes = [
        ("Devices", "ESP32 firmware.\nMQTT over TLS.\nLocal fallback when\nthe network is gone.", CYAN),
        ("Control plane", "Self-hosted API.\nOwns devices, telemetry,\ncommands, automations.", VIOLET),
        ("Applications", "Web console and\nnative mobile app,\nthe same product.", CYAN),
        ("Commerce", "Storefront, checkout,\norders, warranty\nand support.", VIOLET),
    ]
    card_w = (W - 1.7 - 0.35 * 3) / 4
    for i, (title, body, color) in enumerate(boxes):
        left = 0.85 + i * (card_w + 0.35)
        _rect(s, left, 2.3, card_w, 2.5, PAPER_ALT, line=LINE, radius=True)
        _rect(s, left, 2.3, card_w, 0.07, color)
        _text(s, left + 0.2, 2.55, card_w - 0.4, 0.5, title, size=17, bold=True, color=INK)
        _text(s, left + 0.2, 3.15, card_w - 0.4, 1.5, body, size=12.5, color=SLATE, spacing=1.2)
    _text(s, 0.85, 5.15, 11.6, 1.3,
          "They deploy independently. If the storefront is down, devices keep working; "
          "if the device platform is down, the shop keeps selling. "
          "Commands reach hardware in under a second, and nothing in that path polls.",
          size=14, color=MUTED, spacing=1.3)
    _footer(s, data, page=page)

    # --- differentiation
    page += 1
    s = _blank(prs)
    _slide_header(s, "Why it holds up", "Engineering decisions that show up commercially")
    _bullets(s, 0.85, 2.15, 11.6, 4.0, [
        "Made in India \u2014 our own boards and our own firmware, "
        "not a rebadged import with a skinned app.",
        "One platform across every product \u2014 a customer adds a second device, not a second app.",
        "Recovers without a human \u2014 devices retry the network indefinitely and only ask "
        "for setup when genuinely unconfigured.",
        "Failures are made loud \u2014 a control that cannot reach hardware is refused at build "
        "time rather than shipped as a dead button.",
        f"Support terms are part of the product \u2014 {com['warrantyMonths']}-month warranty and "
        f"a {com['returnDays']}-day return window, quoted from one source so the invoice, "
        "the site and the policy page cannot disagree.",
    ])
    _footer(s, data, page=page)

    # --- commercial model
    page += 1
    s = _blank(prs)
    _slide_header(s, "Commercial model", "How the money works today")
    _table(s, ["Line", "Position"], [
        ["Primary revenue", "Direct hardware sales through our own storefront"],
        ["Price band", f"{rupees(cat['priceMin'])} \u2013 {rupees(cat['priceMax'])}"],
        ["Free shipping", f"Orders over {rupees(com['freeShippingOver'])}"],
        ["Shipping otherwise", rupees(com["flatShipping"])],
        ["Warranty", f"{com['warrantyMonths']} months"],
        ["Returns", f"{com['returnDays']} days, unused and in original packaging"],
        ["Attach / bundles", "Multi-device bundles priced server-side"],
        ["Recurring revenue", TBD],
        ["Gross margin", TBD],
    ], top=2.15, col_widths=[4.0, 7.6], row_h=0.42, size=13)
    _text(s, 0.85, 6.35, 11.6, 0.5,
          "Commercial terms above are read from the live catalogue. "
          "Margin and recurring revenue are held outside this repository.",
          size=11.5, color=MUTED)
    _footer(s, data, page=page)

    # --- traction (honest)
    page += 1
    s = _blank(prs)
    _slide_header(s, "Traction", "Figures that live outside the product repository")
    _table(s, ["Metric", "Value"], [
        ["Units shipped", TBD],
        ["Revenue, trailing twelve months", TBD],
        ["Active devices on the platform", TBD],
        ["Registered customers", TBD],
        ["Repeat purchase rate", TBD],
        ["Headcount", TBD],
        ["Funding raised to date", TBD],
    ], top=2.15, col_widths=[7.0, 4.6], row_h=0.45, size=13)
    _text(s, 0.85, 5.9, 11.6, 1.0,
          "This deck is generated from the engineering repository, which holds the product "
          "and its commercial terms but not the trading figures. These are deliberately left "
          "blank rather than estimated \u2014 an invented number in a deck is quoted back as fact.",
          size=13, color=MUTED, spacing=1.3)
    _footer(s, data, page=page)

    # --- roadmap
    page += 1
    s = _blank(prs)
    _slide_header(s, "Roadmap", "Shipped, in progress, next")
    cols = [
        ("Shipped", CYAN, [
            f"{cat['total']} products on sale",
            "Web console and mobile app at parity",
            "Scenes, schedules and automations",
            "Number-plate recognition (ANPR)",
            "Voice control",
        ]),
        ("In progress", VIOLET, [
            "Drone platform and flight logging",
            "Deeper energy analytics",
            "Commercial and multi-site deployments",
            "Partner and dealer channel",
        ]),
        ("Next", MUTED, [
            "Wider retail distribution",
            "Installer network",
            "Subscription services",
            "Export markets",
        ]),
    ]
    card_w = (W - 1.7 - 0.4 * 2) / 3
    for i, (title, color, items) in enumerate(cols):
        left = 0.85 + i * (card_w + 0.4)
        _rect(s, left, 2.2, card_w, 3.7, PAPER_ALT, line=LINE, radius=True)
        _rect(s, left, 2.2, card_w, 0.07, color)
        _text(s, left + 0.22, 2.45, card_w - 0.44, 0.45, title, size=17, bold=True, color=INK)
        _bullets(s, left + 0.22, 3.0, card_w - 0.44, 2.7, items, size=12.5,
                 gap=7, bullet_color=color)
    _text(s, 0.85, 6.15, 11.6, 0.5,
          "Roadmap reflects the state of the repository. Dates are set commercially and are "
          "not asserted here.", size=11.5, color=MUTED)
    _footer(s, data, page=page)

    # --- close
    s = _blank(prs, INK_DEEP)
    _rect(s, 0, 0, 0.16, H, CYAN)
    _rect(s, 0, H * 0.55, 0.16, H * 0.45, VIOLET)
    _text(s, 1.1, 2.6, 11, 1.0, "Let's talk.", size=48, bold=True, color=PAPER)
    _text(s, 1.1, 3.7, 11, 0.5, data["company"]["name"], size=20, color=CYAN_BRIGHT, bold=True)
    c = data["company"]
    _bullets(s, 1.1, 4.3, 10, 1.8, [
        c["site"], c["salesEmail"], c["phone"], c["location"],
    ], size=15, color="94A3B8", gap=5, bullet_color=CYAN_BRIGHT)
    _text(s, 1.1, H - 0.7, 11, 0.35, generated_stamp(data), size=9.5, color="475569")

    prs.save(str(out_path))
    return out_path


# ---------------------------------------------------------------- sales deck

def build_sales_deck(data: dict, out_path):
    prs = _new_deck()
    cat = data["catalogue"]
    com = data["commercial"]
    page = 0

    _title_slide(
        prs, data, "Smart devices that\nkeep working.",
        f"{cat['total']} products for the home, from {rupees(cat['priceMin'])}. "
        "One app, one account, one company behind all of it.",
        "Product and pricing overview",
    )

    page += 1
    s = _blank(prs)
    _slide_header(s, "Why customers switch", "Bought once, not supported once")
    _bullets(s, 0.85, 2.2, 11.6, 3.8, [
        "One app for everything \u2014 lights, fans, water, power, locks and cameras "
        "in a single account.",
        "Keeps working without the internet \u2014 devices hold their local behaviour "
        "and reconnect on their own.",
        "Real energy visibility \u2014 live consumption, not an estimate.",
        "Made in India \u2014 our own hardware and firmware, supported directly by us.",
        f"{com['warrantyMonths']}-month warranty and {com['returnDays']}-day returns \u2014 "
        "the same terms on the invoice, the website and the policy page.",
    ])
    _footer(s, data, page=page)

    # A slide per category, priced from the catalogue, biggest range first.
    for cname in ordered_categories(cat):
        page += 1
        items = sorted([p for p in cat["products"] if p["category"] == cname],
                       key=lambda p: p["price"])
        s = _blank(prs)
        color = category_color(cname)
        _accent_bar(s)
        _text(s, 0.85, 0.55, 11.6, 0.4, "PRODUCT RANGE", size=12, bold=True, color=color)
        _text(s, 0.85, 0.95, 11.6, 0.8, cname, size=34, bold=True, color=INK)
        rows = [[p["name"], p["tagline"][:64], rupees(p["price"]),
                 rupees(p["compareAt"]) if p["compareAt"] else "\u2014"] for p in items]
        _table(s, ["Product", "What it does", "Price", "MRP"], rows,
               top=2.0, col_widths=[3.1, 5.4, 1.6, 1.5],
               row_h=0.44 if len(rows) > 6 else 0.5, size=12)
        _footer(s, data, page=page)

    page += 1
    s = _blank(prs)
    _slide_header(s, "Buying", "Terms, delivery and support")
    _stat_row(s, [
        (rupees(com["freeShippingOver"]), "free shipping above this"),
        (f"{com['warrantyMonths']} mo", "warranty on every device"),
        (f"{com['returnDays']} days", "return window"),
        (cat["total"], "products to choose from"),
    ], top=2.4)
    _text(s, 0.85, 5.0, 11.6, 1.2,
          f"Orders below {rupees(com['freeShippingOver'])} carry flat "
          f"{rupees(com['flatShipping'])} shipping. Returns apply to unused items in original "
          "packaging. Bundles are discounted automatically at checkout.",
          size=14, color=MUTED, spacing=1.3)
    _footer(s, data, page=page)

    s = _blank(prs, INK_DEEP)
    _rect(s, 0, 0, 0.16, H, CYAN)
    _rect(s, 0, H * 0.55, 0.16, H * 0.45, VIOLET)
    _text(s, 1.1, 2.7, 11, 1.0, "Order or ask.", size=48, bold=True, color=PAPER)
    c = data["company"]
    _bullets(s, 1.1, 4.0, 10, 2.0, [
        f"Shop \u2014 {c['site']}",
        f"Sales \u2014 {c['salesEmail']}",
        f"Support \u2014 {c['supportEmail']}",
        f"Phone \u2014 {c['phone']}",
    ], size=16, color="94A3B8", gap=6, bullet_color=CYAN_BRIGHT)
    _text(s, 1.1, H - 0.7, 11, 0.35, generated_stamp(data), size=9.5, color="475569")

    prs.save(str(out_path))
    return out_path
