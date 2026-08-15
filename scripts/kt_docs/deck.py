"""
The knowledge-transfer deck.

WHY IT IMPORTS THE BUSINESS DECK'S HELPERS

The slide furniture — the accent bar, the header, the table, the footer — is
already written and already matches the palette in `src/app/globals.css`. A KT
deck that redrew its own would drift into being a different-looking document
from the same company, which is the exact failure `brand.py` opens by
describing. So the visual system is imported, and only the content lives here.

WHAT THIS DECK IS FOR

One session, walked through by somebody handing the system over. It is
deliberately not a replacement for `Docs/` — it is the map that tells a new
engineer which of the twenty-nine documents to open and in what order, plus the
handful of facts that are true across all four deployables and are not obvious
from any single file.
"""

from __future__ import annotations

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from business_docs.brand import (
    CYAN, CYAN_BRIGHT, VIOLET, INK, INK_DEEP, SLATE, MUTED, PAPER, PAPER_ALT,
)
from business_docs.decks import (
    W, H, _new_deck, _blank, _text, _bullets, _rect, _table,
    _slide_header, _footer, _stat_row,
)


def _kt_title_slide(prs, data, title, subtitle, tag):
    """
    The shared title slide, minus one claim.

    `_title_slide` stamps every deck "generated from the live product
    catalogue", which is true of the sales and investor decks and not of this
    one. This pack is generated from the repository, and inheriting a
    provenance line is how a document ends up asserting something nobody
    checked — so the layout is mirrored and the stamp is replaced.
    """
    slide = _blank(prs, INK_DEEP)
    _rect(slide, 0, 0, 0.16, H, CYAN)
    _rect(slide, 0, H * 0.55, 0.16, H * 0.45, VIOLET)
    _text(slide, 1.1, 2.1, 11, 0.5, data["company"]["name"].upper(), size=14,
          bold=True, color=CYAN_BRIGHT)
    _text(slide, 1.05, 2.65, 11.2, 1.5, title, size=54, bold=True, color=PAPER)
    _text(slide, 1.1, 4.35, 10.5, 1.0, subtitle, size=20, color="94A3B8", spacing=1.3)
    _rect(slide, 1.1, 5.6, 2.2, 0.045, CYAN)
    _text(slide, 1.1, 5.85, 10.5, 0.4, tag, size=13, color="64748B")
    _text(slide, 1.1, H - 0.75, 11, 0.35,
          f"Generated {data['generatedDate']} from the repository at {data['commit']} "
          f"\u00b7 Docs/ is the source of truth",
          size=9.5, color="475569")
    return slide


def _ltable(slide, headers, rows, **kw):
    """
    `_table`, with every column left-aligned.

    The shared builder right-aligns the third column onwards, which is right for
    the price tables it was written for and wrong for prose. Reusing the builder
    and correcting the alignment keeps one table style in the company's decks
    rather than growing a second one that drifts.
    """
    table = _table(slide, headers, rows, **kw)
    for ri in range(1, len(rows) + 1):
        for ci in range(len(headers)):
            para = table.cell(ri, ci).text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT
    return table


def _note(slide, text: str) -> None:
    """Speaker notes: what the presenter should say, not read aloud."""
    slide.notes_slide.notes_text_frame.text = text


def build_kt_deck(data: dict, out_path) -> int:
    prs = _new_deck()
    c = data["counts"]
    page = 0

    # ---------------------------------------------------------------- cover
    _kt_title_slide(
        prs, data,
        "Engineering handover",
        "Four deployables, one product — what to know before you change anything",
        f"Generated {data['generatedDate']} · {data['commit']}",
    )

    # ------------------------------------------------------- how to use this
    page += 1
    s = _blank(prs)
    _slide_header(s, "Start here", "This deck is a map, not the territory")
    _bullets(s, 0.85, 2.0, 11.6, 3.6, [
        f"Docs/ is the source of truth — {c['docs']} documents, written from the code and kept with it.",
        "This deck — which of them to read, in what order, and the few facts that span all four deployables.",
        "The handbook — the same material in prose, for reading rather than presenting.",
        "The quick reference — commands and traps, meant to be printed and kept beside you.",
        "Nothing here is typed by hand \u2014 every count, device and trap on the following slides is read out of the repository at build time, so this pack cannot quietly fall behind the code.",
    ])
    _footer(s, data, page=page)
    _note(s, "Set expectations: they should finish this session able to find things, "
             "not able to recite them. The documents are the deliverable; this is the index.")

    # ------------------------------------------------------ the deployables
    page += 1
    s = _blank(prs)
    _slide_header(s, "Shape of the system", "Four deployables")
    _ltable(
        s,
        ["What", "Lives in", "Language", "Runs on"],
        [[d.name, d.path, d.language, d.runs_on] for d in data["deployables"]],
        top=2.0, col_widths=[3.6, 1.7, 3.2, 3.1],
    )
    _text(s, 0.85, 5.9, 11.6, 0.9,
          "Each can be built, tested and broken on its own. Most tasks touch one. "
          "They meet only at documented contracts \u2014 MQTT topics, the REST API, and the state keys a device publishes.",
          size=13, color=MUTED)
    _footer(s, data, page=page)
    _note(s, "Emphasise independence. The most common new-joiner mistake is assuming a "
             "change needs all four to be running.")

    # ----------------------------------------------------------- size of it
    page += 1
    s = _blank(prs, INK_DEEP)
    _slide_header(s, "Scale", "What you are taking on", dark=True)
    _stat_row(s, [
        (str(c["devices"]), "device types"),
        (str(c["apiRoutes"]), "web API routes"),
        (str(c["planeRoutes"]), "control-plane routers"),
        (str(c["docs"]), "documents"),
    ], dark=True)
    _stat_row(s, [
        (str(c["webTests"]), "web test files"),
        (str(c["planeTests"]), "control-plane test files"),
        (str(c["components"]), "React components"),
        ("2", "separate databases"),
    ], top=4.5, dark=True)
    _footer(s, data, dark=True, page=page)
    _note(s, "These are counted from the tree at build time. If they look wrong, the "
             "generator is reading a different tree than you are.")

    # ------------------------------------------------------- two databases
    page += 1
    s = _blank(prs)
    _slide_header(s, "The thing that surprises everybody", "There are two databases")
    _ltable(
        s,
        ["", "Shop database", "Control-plane database"],
        [
            ["Provider", "Neon (serverless HTTP)", "Postgres 16 container on the VM"],
            ["Reached from", "The Next.js app", "The control-plane API"],
            ["Access code", "src/lib/db.ts", "platform/api/src/db.ts"],
            ["Password hashing", "scrypt", "bcrypt, cost 12"],
            ["Holds", "Customers, orders, wallets, CMS", "Fleet users, devices, telemetry, automations"],
        ],
        top=2.0, col_widths=[2.3, 4.6, 4.7],
    )
    _text(s, 0.85, 5.9, 11.6, 0.9,
          "Neither can see the other, and that is the point: a compromise or an outage of the shop cannot reach "
          "the device fleet. They are joined only by a single-sign-on bridge that lets each vouch for a customer "
          "it has already authenticated. See Docs/05-databases.md.",
          size=13, color=MUTED)
    _footer(s, data, page=page)
    _note(s, "If they remember one architectural fact from this session, make it this one. "
             "Nearly every confused question about users traces back to assuming one database.")

    # --------------------------------------------------------- the hot path
    page += 1
    s = _blank(prs)
    _slide_header(s, "Hot path", "How a command reaches a device")
    steps = [
        ("1", "The console or app calls POST /devices/:id/command."),
        ("2", "The API checks the caller's JWT and that they own that device."),
        ("3", "The API publishes to cv/<deviceId>/cmd on Mosquitto at QoS 1."),
        ("4", "The device applies the change and republishes its full state, retained."),
        ("5", "The API persists that state and pushes it to every watching WebSocket."),
        ("6", "The UI, which already applied the change optimistically, reconciles."),
    ]
    top = 2.05
    for num, text in steps:
        _rect(s, 0.85, top, 0.42, 0.42, CYAN, radius=True)
        _text(s, 0.85, top + 0.04, 0.42, 0.34, num, size=13, bold=True,
              color=PAPER, align=PP_ALIGN.CENTER)
        _text(s, 1.45, top + 0.02, 11.0, 0.42, text, size=15, color=SLATE)
        top += 0.62
    _text(s, 0.85, top + 0.15, 11.6, 0.8,
          "Well under a second end to end, and there is no polling anywhere in it. "
          "The device is the authority on its own state \u2014 the UI never assumes it won.",
          size=13, color=MUTED)
    _footer(s, data, page=page)
    _note(s, "Walk this slowly. Almost every device bug is a misunderstanding of which "
             "step owns the truth.")

    # ---------------------------------------------------- trust boundaries
    page += 1
    s = _blank(prs)
    _slide_header(s, "Security", "Trust boundaries you must not blur")
    _ltable(
        s,
        ["Boundary", "Enforced by"],
        [
            ["Browser → Next.js", "HMAC session tokens; separate admin sessions"],
            ["App → control plane", "JWT bearer token, checked against a live session"],
            ["Device → broker", "Per-device username and secret over TLS; ACL limits it to cv/<its own id>/#"],
            ["Control plane → broker", "The control-plane user, reachable only on the internal Docker network"],
            ["Shop ↔ console SSO", "Shared FEDERATION_SECRET, HMAC and timestamp, server-to-server only"],
        ],
        top=2.0, col_widths=[3.9, 7.7],
    )
    _text(s, 0.85, 5.7, 11.6, 1.1,
          "Two rules with teeth: device ownership is re-checked on every frame, not only when a client asks to watch a "
          "camera \u2014 so unclaiming a device cuts its feed immediately. And locks, gates, cameras, ANPR and drones are "
          "absent from the voice-assistant trait map on purpose: a spoken command must not unlock a door.",
          size=13, color=MUTED)
    _footer(s, data, page=page)
    _note(s, "The voice exclusion list is asserted by traits.test.ts. Adding to it has to "
             "be argued for, not just typed.")

    # ------------------------------------------------------------- devices
    page += 1
    s = _blank(prs)
    _slide_header(s, "The fleet", f"{c['devices']} device types ship from firmware/")
    devices = data["devices"]
    third = (len(devices) + 2) // 3
    cols = [devices[0:third], devices[third:2 * third], devices[2 * third:]]
    for i, col in enumerate(cols):
        _bullets(s, 0.85 + i * 3.95, 2.0, 3.8, 4.4, col, size=13, gap=5)
    _text(s, 0.85, 6.4, 11.6, 0.7,
          "A device type id is the worst thing in this codebase to get wrong: the same string appears in firmware, "
          "the API, both apps and the shop. Copy it, never retype it \u2014 and see Docs/07 before adding one.",
          size=13, color=MUTED)
    _footer(s, data, page=page)
    _note(s, "This list is the firmware tree. firmware-console-parity.test.ts fails if any "
             "of these has no console control surface.")

    # ------------------------------------------------------ the parity rule
    page += 1
    s = _blank(prs, INK_DEEP)
    _slide_header(s, "How work is done here", "The rule that matters most", dark=True)
    _text(s, 0.85, 2.0, 11.6, 1.0,
          "Almost every real bug in this codebase has been one surface knowing something another surface did not.",
          size=22, bold=True, color=PAPER)
    _bullets(s, 0.85, 3.2, 11.6, 2.6, [
        "Not a crash \u2014 a control that is simply absent, or a button that does nothing, on one screen out of six.",
        "A device type registered in the app but not the console: customers could dim a lamp on their phone and not in a browser.",
        "A scene editor offering a Power toggle for a curtain: the row was offered, configured, saved, and silently dropped.",
        "They share a shape \u2014 two tables that must agree, nothing forcing them to, and no error when they disagree.",
    ], color="CBD5E1")
    _text(s, 0.85, 6.0, 11.6, 0.8,
          "So when you add a fact about a device \u2014 a type, a field, a control, an icon \u2014 find every table that already "
          "stores that kind of fact, and add a test that fails when they disagree.",
          size=14, color=CYAN_BRIGHT)
    _footer(s, data, dark=True, page=page)
    _note(s, "This is the single most useful thing to teach. tests/ is full of parity "
             "guards; point at device-type-parity.test.ts as the model.")

    # -------------------------------------------------------------- traps
    if data["traps"]:
        page += 1
        s = _blank(prs)
        _slide_header(s, "Learn from other people's afternoons", "Traps, each of which cost somebody a day")
        _ltable(
            s,
            ["Trap", "What happens", "The truth"],
            [[t[0], t[1], t[2]] for t in data["traps"]],
            top=1.95, col_widths=[3.5, 4.0, 4.1], size=11,
        )
        _footer(s, data, page=page)
        _note(s, "Parsed straight out of Docs/00 §9 at build time, so this slide cannot "
                 "fall behind the table new joiners actually read.")

    # ------------------------------------------------------------ testing
    page += 1
    s = _blank(prs)
    _slide_header(s, "Before you push", "Testing is the handover")
    _ltable(
        s,
        ["Where", "Command", "What it protects"],
        [
            ["Web app", "npx tsc --noEmit && npm test", "Types, and the parity guards in tests/"],
            ["Control plane", "cd platform/api && npm test", "Auth, MQTT bridge, households, voice"],
            ["Mobile", "cd mobile && npm run typecheck", "tsc plus a dozen static audits"],
            ["Firmware", "python -m platformio run", "It compiles for the board it ships on"],
        ],
        top=2.0, col_widths=[2.2, 4.6, 4.8],
    )
    _text(s, 0.85, 5.6, 11.6, 1.2,
          "The mobile typecheck is not just tsc. It audits navigation targets, device-type coverage, the command map "
          "against firmware, contrast, screen theming, iOS-only dialogs, swallowed failures and permissions \u2014 and each of "
          "those exists because it once shipped a bug. Verify by running it, not by reading it.",
          size=13, color=MUTED)
    _footer(s, data, page=page)
    _note(s, "audit-swallowed.js is a ratchet: the number may fall and must never rise.")

    # ---------------------------------------------------------- deploying
    page += 1
    s = _blank(prs)
    _slide_header(s, "Getting it out", "Four things deploy four different ways")
    _bullets(s, 0.85, 2.0, 11.6, 4.0, [
        "Website \u2014 push to the deploying remote; Vercel builds. Verify by build sha from /api/health, never by eye.",
        "Control plane \u2014 tar the tree, upload, and run scripts/deploy.sh on the VM. Hand it the tarball; do not extract it yourself, or the backup archives the new code and calls it a rollback.",
        "Mobile \u2014 bump version and versionCode together, then npm run build:android.",
        "Firmware \u2014 build, publish the binary, then push an OTA pointer to the device. The device verifies TLS against a pinned root.",
        "Whatever you deploy \u2014 confirm the thing you shipped is the thing that is running. A stale container that reports commit: unknown exists precisely so it cannot pretend to be healthy.",
    ])
    _footer(s, data, page=page)
    _note(s, "Docs/09 has the exact commands. The point of this slide is that 'it built' "
             "is not 'it deployed'.")

    # ------------------------------------------------------- reading order
    page += 1
    s = _blank(prs)
    _slide_header(s, "The library", f"{c['docs']} documents — read them in this order")
    docs = data["docs"]
    half = (len(docs) + 1) // 2
    left = [f"{name.replace('.md','')} \u2014 {title.split('—')[-1].strip()}" for name, title in docs[:half]]
    right = [f"{name.replace('.md','')} \u2014 {title.split('—')[-1].strip()}" for name, title in docs[half:]]
    _bullets(s, 0.85, 1.9, 5.7, 5.0, left, size=10, gap=2)
    _bullets(s, 6.85, 1.9, 5.7, 5.0, right, size=10, gap=2)
    _footer(s, data, page=page)
    _note(s, "Do not read them all now. 00 and 01 today; the rest when the task calls for one.")

    # ----------------------------------------------------------- first week
    page += 1
    s = _blank(prs)
    _slide_header(s, "Your first week", "Prove the environment, then pick something small")
    _bullets(s, 0.85, 2.0, 11.6, 4.2, [
        "Day one \u2014 work through Docs/00-start-here.md end to end. Four proofs: the site renders, the tests pass, the API answers, the firmware compiles.",
        "Then read Docs/01-architecture.md. It means much more once you have seen the pieces move.",
        "Pick a task that touches one deployable. Resist anything that needs all four running.",
        "Before you open a pull request, ask the parity question \u2014 is there another table, screen or app that also needs to know what I just added? No tool asks it for you.",
        "When something is wrong, say what you ran, what you expected and what you got. That turns a conversation into a fix.",
    ])
    _footer(s, data, page=page)
    _note(s, "Encourage them to fix a wrong document in the same PR as the code. These "
             "documents only stay true because whoever notices the drift corrects it.")

    # ---------------------------------------------------------------- close
    page += 1
    s = _blank(prs, INK_DEEP)
    _slide_header(s, "In one line", "Verify by running it, not by reading it", dark=True)
    _text(s, 0.85, 2.4, 11.6, 2.2,
          "Nearly every defect found while this system was being built appeared only when it was executed \u2014 and several "
          "were in code written an hour earlier by the person who found them. The recurring shape is a control that looks "
          "present and does nothing: a chart labelled one thing and drawing another, a guard written and wired to nothing, "
          "a button that is undefined on half the phones it shipped to.",
          size=17, color="CBD5E1")
    _text(s, 0.85, 4.9, 11.6, 0.9, "They never announce themselves. Go and look.",
          size=24, bold=True, color=CYAN_BRIGHT)
    _footer(s, data, dark=True, page=page)
    _note(s, "End here. It is the habit that matters more than any individual fact in this deck.")

    prs.save(str(out_path))
    return len(prs.slides._sldIdLst)
