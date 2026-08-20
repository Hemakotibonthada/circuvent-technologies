#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
generate_docs.py — Multi-format documentation builder for Career.circuvent.

Reads the numbered Markdown sources in this directory and produces:

    Architecture_Guide.md       aggregated single-page master reference
    Architecture_Guide.docx     styled Word document
    Architecture_Guide.pdf      cover page, table of contents, vector diagrams
    Architecture_Overview.pptx  14-slide presentation deck

Self-contained: no network access, no build step. Run it from anywhere.

    python generate_docs.py

Dependencies: python-docx, python-pptx, reportlab, markdown
"""

from __future__ import annotations

import os
import re
import sys
import datetime
from dataclasses import dataclass, field
from typing import Iterable

# ───────────────────────────────────────────────────────────── constants ──

HERE = os.path.dirname(os.path.abspath(__file__))

SOURCES = [
    ("01_SYSTEM_OVERVIEW.md", "System Overview"),
    ("02_DATABASE_AND_DATA_MODELS.md", "Database & Data Models"),
    ("03_INTEGRATIONS_AND_ECOSYSTEM.md", "Integrations & Ecosystem"),
    ("04_MAINTENANCE_AND_OPERATIONS.md", "Maintenance & Operations"),
    ("05_AREAS_OF_ENHANCEMENT.md", "Areas of Enhancement"),
    ("06_ARCHITECTURE_DIAGRAMS.md", "Architecture Diagram Atlas"),
]

TITLE = "Circuvent Technologies"
SUBTITLE = "Architecture & Technical Audit"
ORG = "Circuvent Technologies"
BUILT = datetime.date.today().isoformat()

# Brand palette (hex, no leading #)
INK = "1F2933"
MUTED = "5A6875"
ACCENT = "1D4ED8"
ACCENT_DK = "1E3A8A"
GOOD = "15803D"
WARN = "B45309"
BAD = "B91C1C"
RULE = "D5DBE1"
CODE_BG = "F4F6F8"
BAND = "0F172A"

# Glyphs that no installed font renders, mapped per-context. Verified against
# the actual cmap tables of Consolas and Calibri with fontTools, not guessed.
#
#   mono  — used inside code blocks. MUST be a single character, because these
#           diagrams are column-aligned and a two-character substitution shears
#           every box below it.
#   prose — used in body text, tables, headings and slides, where a word is
#           clearer and width does not matter.
#
# Consolas has the full box-drawing and block-element range, so those survive
# in code. Calibri does not, so they are flattened in prose.
GLYPHS = {
    # (character): (mono replacement, prose replacement)
    "\u2705": ("Y", "YES"),          # ✅
    "\u274c": ("N", "NO"),           # ❌
    "\u26a0": ("!", "!"),            # ⚠
    "\u2713": ("+", "yes"),          # ✓
    "\u2717": ("x", "no"),           # ✗
    "\u2605": ("*", "*"),            # ★
    "\u2606": (".", "-"),            # ☆
    "\u25b6": (">", ">"),            # ▶
    "\u25c0": ("<", "<"),            # ◀
    "\u1f512": ("#", "[encrypted]"),  # 🔒
    "\u1f534": ("!", "[CRITICAL]"),  # 🔴
    "\U0001f7e0": ("!", "[HIGH]"),   # 🟠
    "\U0001f7e1": ("~", "[MEDIUM]"),  # 🟡
    "\U0001f7e2": (".", "[LOW]"),    # 🟢
    "\u2699": ("*", ""),             # ⚙
    "\U0001f310": ("*", ""),         # 🌐
    "\U0001f4e6": ("*", "[pkg]"),    # 📦
    "\ufe0f": ("", ""),              # variation selector
    "\U0001f535": (".", "[INFO]"),  # 🔵
    "\u2b50": ("*", "*"),           # ⭐
    "\u2716": ("x", "x"),           # ✖
    "\u27e8": ("<", "<"),           # ⟨
    "\u27e9": (">", ">"),           # ⟩
    "\u2208": ("in", "in"),         # ∈
}

# Present in Consolas but absent from Calibri: flattened in prose only.
PROSE_ONLY = {
    "\u2550": "-", "\u2551": "|", "\u2588": "#", "\u2591": ":",
    "\u25bc": "v", "\u25b2": "^", "\u2570": "'", "\u256f": "'",
    "\u2554": "+", "\u2557": "+", "\u255a": "+", "\u255d": "+",
    "\u2560": "+", "\u2563": "+", "\u2564": "+",
}


def strip_glyphs(text: str, mono: bool = False) -> str:
    """Replace characters the target font cannot draw.

    `mono=True` keeps box-drawing intact (Consolas has it) and uses
    single-character substitutions so diagram columns stay aligned.
    """
    idx = 0 if mono else 1
    for ch, pair in GLYPHS.items():
        if ch in text:
            text = text.replace(ch, pair[idx])
    if not mono:
        for ch, rep in PROSE_ONLY.items():
            if ch in text:
                text = text.replace(ch, rep)
    return text


# ─────────────────────────────────────────────────────────── md parsing ──

@dataclass
class Block:
    kind: str                       # heading|para|code|table|list|quote|hr
    text: str = ""
    level: int = 0
    lang: str = ""
    lines: list = field(default_factory=list)
    rows: list = field(default_factory=list)
    ordered: bool = False


def _split_row(line: str) -> list:
    line = line.strip()
    if line.startswith("|"):
        line = line[1:]
    if line.endswith("|"):
        line = line[:-1]
    return [c.strip() for c in line.split("|")]


def parse_markdown(text: str) -> list:
    """Line-based Markdown parser covering the subset used by these documents."""
    out: list = []
    lines = text.replace("\r\n", "\n").split("\n")
    i = 0
    n = len(lines)

    while i < n:
        line = lines[i]
        stripped = line.strip()

        # fenced code
        if stripped.startswith("```"):
            lang = stripped[3:].strip()
            i += 1
            buf = []
            while i < n and not lines[i].strip().startswith("```"):
                buf.append(lines[i])
                i += 1
            i += 1
            while buf and not buf[-1].strip():
                buf.pop()
            out.append(Block("code", lang=lang, lines=buf))
            continue

        # blank
        if not stripped:
            i += 1
            continue

        # horizontal rule
        if re.fullmatch(r"(-{3,}|\*{3,}|_{3,})", stripped):
            out.append(Block("hr"))
            i += 1
            continue

        # heading
        m = re.match(r"^(#{1,6})\s+(.*)$", stripped)
        if m:
            out.append(Block("heading", text=m.group(2).strip(), level=len(m.group(1))))
            i += 1
            continue

        # table: a header row followed by a separator row
        if stripped.startswith("|") and i + 1 < n and re.match(
            r"^\s*\|?[\s:\-|]+\|[\s:\-|]*$", lines[i + 1]
        ) and "-" in lines[i + 1]:
            header = _split_row(lines[i])
            i += 2
            rows = []
            while i < n and lines[i].strip().startswith("|"):
                rows.append(_split_row(lines[i]))
                i += 1
            width = len(header)
            norm = []
            for r in rows:
                r = (r + [""] * width)[:width]
                norm.append(r)
            out.append(Block("table", rows=[header] + norm))
            continue

        # blockquote
        if stripped.startswith(">"):
            buf = []
            while i < n and lines[i].strip().startswith(">"):
                buf.append(re.sub(r"^\s*>\s?", "", lines[i]))
                i += 1
            merged = []
            cur = []
            for b in buf:
                if b.strip():
                    cur.append(b.strip())
                else:
                    if cur:
                        merged.append(" ".join(cur))
                        cur = []
            if cur:
                merged.append(" ".join(cur))
            out.append(Block("quote", lines=merged))
            continue

        # list
        if re.match(r"^\s*([-*+]|\d+\.)\s+", line):
            ordered = bool(re.match(r"^\s*\d+\.\s+", line))
            items = []
            while i < n and re.match(r"^\s*([-*+]|\d+\.)\s+", lines[i]):
                raw = lines[i]
                indent = len(raw) - len(raw.lstrip())
                body = re.sub(r"^\s*([-*+]|\d+\.)\s+", "", raw).strip()
                i += 1
                # continuation lines
                while (i < n and lines[i].strip()
                       and not re.match(r"^\s*([-*+]|\d+\.)\s+", lines[i])
                       and not lines[i].strip().startswith(("#", "|", "```", ">"))
                       and (len(lines[i]) - len(lines[i].lstrip())) > indent):
                    body += " " + lines[i].strip()
                    i += 1
                items.append((min(indent // 2, 2), body))
            out.append(Block("list", rows=items, ordered=ordered))
            continue

        # paragraph
        buf = []
        while (i < n and lines[i].strip()
               and not lines[i].strip().startswith(("#", "```", ">", "|"))
               and not re.match(r"^\s*([-*+]|\d+\.)\s+", lines[i])
               and not re.fullmatch(r"(-{3,}|\*{3,}|_{3,})", lines[i].strip())):
            buf.append(lines[i].strip())
            i += 1
        if buf:
            out.append(Block("para", text=" ".join(buf)))
        else:
            i += 1

    return out


# inline: -> list of (text, bold, italic, code, href)
_INLINE = re.compile(
    r"(?P<link>\[[^\]]+\]\([^)]+\))"
    r"|(?P<code>`[^`]+`)"
    r"|(?P<bold>\*\*[^*]+\*\*)"
    r"|(?P<italic>(?<!\*)\*[^*\n]+\*(?!\*))"
)


def inline_tokens(text: str):
    pos = 0
    for m in _INLINE.finditer(text):
        if m.start() > pos:
            yield (text[pos:m.start()], False, False, False, None)
        raw = m.group(0)
        if m.group("code"):
            yield (raw[1:-1], False, False, True, None)
        elif m.group("link"):
            lm = re.match(r"\[([^\]]+)\]\(([^)]+)\)", raw)
            label = lm.group(1).strip()
            if label.startswith("`") and label.endswith("`"):
                label = label[1:-1]
            yield (label, False, False, False, lm.group(2))
        elif m.group("bold"):
            # recurse: bold spans may wrap code, links or italics
            for t, _b, i, c, h in inline_tokens(raw[2:-2]):
                yield (t, True, i, c, h)
        else:
            for t, b, _i, c, h in inline_tokens(raw[1:-1]):
                yield (t, b, True, c, h)
        pos = m.end()
    if pos < len(text):
        yield (text[pos:], False, False, False, None)


# ────────────────────────────────────────────── 1 · master markdown ──

def build_master_md() -> str:
    parts = [
        f"# {TITLE} — {SUBTITLE}\n",
        f"> **Organisation:** {ORG}  ",
        f"> **Generated:** {BUILT}  ",
        "> **Scope:** full technical audit and architecture reverse-engineering.\n",
        "\nThis is the aggregated master reference. The same content is maintained "
        "as five focused documents in this directory; edit those, then re-run "
        "`generate_docs.py` to rebuild this file and the Word, PDF and PowerPoint "
        "deliverables.\n",
        "\n---\n",
        "\n## Contents\n",
    ]
    for idx, (_, label) in enumerate(SOURCES, start=1):
        anchor = "part-%d-%s" % (idx, re.sub(r"[^a-z0-9]+", "-", label.lower()).strip("-"))
        parts.append(f"{idx}. [Part {idx} · {label}](#{anchor})")
    parts.append("\n---\n")

    for idx, (fname, label) in enumerate(SOURCES, start=1):
        with open(os.path.join(HERE, fname), "r", encoding="utf-8") as fh:
            body = fh.read()
        # cross-document links become plain text inside a single file
        body = re.sub(r"\[`?([^\]`]+)`?\]\(\./\d\d_[A-Z_]+\.md\)", r"**\1**", body)
        body = re.sub(r"^#\s+.*$", "", body, count=1, flags=re.M).lstrip("\n")
        parts.append(f"\n<a id=\"part-{idx}-"
                     f"{re.sub(r'[^a-z0-9]+', '-', label.lower()).strip('-')}\"></a>\n")
        parts.append(f"\n# Part {idx} · {label}\n")
        parts.append(body.rstrip())
        parts.append("\n\n---\n")

    parts.append(f"\n*Generated by `generate_docs.py` on {BUILT}.*\n")
    return "\n".join(parts)


# ────────────────────────────────────────────────────── 2 · docx ──

def build_docx(path: str) -> None:
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.enum.section import WD_SECTION
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()

    sec = doc.sections[0]
    sec.left_margin = Inches(0.85)
    sec.right_margin = Inches(0.85)
    sec.top_margin = Inches(0.85)
    sec.bottom_margin = Inches(0.8)

    normal = doc.styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = Pt(10.5)
    normal.font.color.rgb = RGBColor.from_string(INK)
    normal.paragraph_format.space_after = Pt(6)
    normal.paragraph_format.line_spacing = 1.15
    rpr = normal.element.get_or_add_rPr()
    rf = rpr.find(qn("w:rFonts"))
    if rf is None:
        rf = OxmlElement("w:rFonts")
        rpr.append(rf)
    rf.set(qn("w:eastAsia"), "Calibri")

    for name, size, color, bold in (
        ("Heading 1", 20, ACCENT_DK, True),
        ("Heading 2", 15, ACCENT_DK, True),
        ("Heading 3", 12.5, INK, True),
        ("Heading 4", 11, MUTED, True),
        ("Heading 5", 10.5, MUTED, True),
        ("Heading 6", 10, MUTED, True),
    ):
        st = doc.styles[name]
        st.font.name = "Calibri"
        st.font.size = Pt(size)
        st.font.bold = bold
        st.font.color.rgb = RGBColor.from_string(color)
        st.paragraph_format.space_before = Pt(14 if name == "Heading 1" else 10)
        st.paragraph_format.space_after = Pt(6)
        st.paragraph_format.keep_with_next = True

    def shade(paragraph, hexcolor):
        pr = paragraph._p.get_or_add_pPr()
        sh = OxmlElement("w:shd")
        sh.set(qn("w:val"), "clear")
        sh.set(qn("w:fill"), hexcolor)
        pr.append(sh)

    def left_bar(paragraph, hexcolor):
        pr = paragraph._p.get_or_add_pPr()
        bd = OxmlElement("w:pBdr")
        lf = OxmlElement("w:left")
        lf.set(qn("w:val"), "single")
        lf.set(qn("w:sz"), "18")
        lf.set(qn("w:space"), "8")
        lf.set(qn("w:color"), hexcolor)
        bd.append(lf)
        pr.append(bd)

    def add_runs(paragraph, text, base_size=10.5, mono=False):
        text = strip_glyphs(text)
        for chunk, bold, italic, code, href in inline_tokens(text):
            if not chunk:
                continue
            run = paragraph.add_run(chunk)
            run.bold = bold
            run.italic = italic
            if code or mono:
                run.font.name = "Consolas"
                run.font.size = Pt(base_size - 1)
                run.font.color.rgb = RGBColor.from_string(BAD if code else INK)
                r = run._element.get_or_add_rPr().get_or_add_rFonts()
                r.set(qn("w:ascii"), "Consolas")
                r.set(qn("w:hAnsi"), "Consolas")
            else:
                run.font.size = Pt(base_size)
            if href:
                run.font.color.rgb = RGBColor.from_string(ACCENT)
                run.underline = True

    # ── cover ──
    for _ in range(5):
        doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(TITLE)
    r.font.size = Pt(34)
    r.font.bold = True
    r.font.color.rgb = RGBColor.from_string(ACCENT_DK)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(SUBTITLE)
    r.font.size = Pt(17)
    r.font.color.rgb = RGBColor.from_string(MUTED)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("─" * 30)
    r.font.color.rgb = RGBColor.from_string(RULE)

    for line, size, color in (
        (ORG, 12, INK),
        ("Website · commerce · IoT platform · firmware · hardware · mobile", 11, MUTED),
        (f"Generated {BUILT}", 10, MUTED),
    ):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run(line)
        r.font.size = Pt(size)
        r.font.color.rgb = RGBColor.from_string(color)

    doc.add_paragraph()
    facts = doc.add_table(rows=0, cols=2)
    facts.style = "Light List Accent 1"
    facts.alignment = WD_TABLE_ALIGNMENT.CENTER
    for k, v in (
        ("Architecture", "Five web products, firmware, hardware and two native apps in one repository"),
        ("Scale", "23,184 files · 1,810 src TypeScript · 428,355 lines · 150 API routes · 108 pages"),
        ("Database", "Neon Postgres, HTTP driver · 10 tables · CREATED AT RUNTIME · no migrations"),
        ("Sessions", "Home-grown HMAC-SHA256 — no JWT library is used anywhere in src/"),
        ("Devices", "17 retail SKUs · ESP32 · MQTT/TLS on an own broker with an own CA"),
        ("Tests", "236 suites · 4,328 tests · 1 failing · 49 seconds"),
        ("CI/CD", "14 steps, 13 hard gates — and 27 of 27 runs failed to start"),
        ("Primary risk", "~27 storage modules are memory-only in production, passkeys among them"),
    ):
        row = facts.add_row().cells
        run = row[0].paragraphs[0].add_run(k)
        run.bold = True
        run.font.size = Pt(9.5)
        run2 = row[1].paragraphs[0].add_run(v)
        run2.font.size = Pt(9.5)

    doc.add_page_break()

    # ── table of contents field ──
    h = doc.add_paragraph()
    r = h.add_run("Contents")
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = RGBColor.from_string(ACCENT_DK)

    p = doc.add_paragraph()
    fld = OxmlElement("w:fldSimple")
    fld.set(qn("w:instr"), r'TOC \o "1-3" \h \z \u')
    inner = OxmlElement("w:r")
    t = OxmlElement("w:t")
    t.text = "Right-click and choose “Update Field” to build the table of contents."
    inner.append(t)
    fld.append(inner)
    p._p.append(fld)

    doc.add_page_break()

    # ── body ──
    for idx, (fname, label) in enumerate(SOURCES, start=1):
        with open(os.path.join(HERE, fname), "r", encoding="utf-8") as fh:
            blocks = parse_markdown(fh.read())

        if idx > 1:
            doc.add_page_break()

        band = doc.add_paragraph()
        shade(band, BAND)
        r = band.add_run(f"  PART {idx}")
        r.font.size = Pt(9)
        r.font.bold = True
        r.font.color.rgb = RGBColor.from_string("FFFFFF")

        first_heading_used = False

        for b in blocks:
            if b.kind == "heading":
                if not first_heading_used and b.level == 1:
                    first_heading_used = True
                    doc.add_heading(strip_glyphs(b.text), level=1)
                    continue
                doc.add_heading(strip_glyphs(b.text), level=min(b.level, 6))

            elif b.kind == "para":
                p = doc.add_paragraph()
                add_runs(p, b.text)

            elif b.kind == "quote":
                for ln in b.lines:
                    p = doc.add_paragraph()
                    p.paragraph_format.left_indent = Inches(0.22)
                    p.paragraph_format.space_before = Pt(4)
                    left_bar(p, ACCENT)
                    shade(p, "EEF3FB")
                    add_runs(p, ln, base_size=10)
                    for run in p.runs:
                        run.italic = True

            elif b.kind == "code":
                is_mermaid = b.lang.lower() == "mermaid"
                if is_mermaid:
                    cap = doc.add_paragraph()
                    cr = cap.add_run("Mermaid diagram — renders in the Markdown edition")
                    cr.font.size = Pt(8)
                    cr.bold = True
                    cr.font.color.rgb = RGBColor.from_string(ACCENT)
                    cap.paragraph_format.space_after = Pt(0)
                widest = max((len(x) for x in b.lines), default=0)
                size = 8 if widest <= 92 else (7 if widest <= 108 else 6)
                for ln in b.lines:
                    p = doc.add_paragraph()
                    p.paragraph_format.space_after = Pt(0)
                    p.paragraph_format.space_before = Pt(0)
                    p.paragraph_format.line_spacing = 1.0
                    shade(p, "EDF2F7" if is_mermaid else CODE_BG)
                    run = p.add_run(strip_glyphs(ln, mono=True) or " ")
                    run.font.name = "Consolas"
                    run.font.size = Pt(size)
                    run.font.color.rgb = RGBColor.from_string(INK)
                    rf2 = run._element.get_or_add_rPr().get_or_add_rFonts()
                    rf2.set(qn("w:ascii"), "Consolas")
                    rf2.set(qn("w:hAnsi"), "Consolas")
                doc.add_paragraph().paragraph_format.space_after = Pt(2)

            elif b.kind == "table":
                header, *rows = b.rows
                has_header = any(c.strip() for c in header)
                if not has_header:
                    rows = [header] + rows
                ncols = len(b.rows[0])
                t = doc.add_table(rows=0, cols=ncols)
                t.style = "Light Grid Accent 1" if has_header else "Light List Accent 1"
                t.autofit = True
                if has_header:
                    cells = t.add_row().cells
                    for c, txt in zip(cells, header):
                        para = c.paragraphs[0]
                        add_runs(para, txt, base_size=9)
                        for run in para.runs:
                            run.bold = True
                for row in rows:
                    cells = t.add_row().cells
                    for c, txt in zip(cells, row):
                        para = c.paragraphs[0]
                        para.paragraph_format.space_after = Pt(2)
                        add_runs(para, txt, base_size=9)
                doc.add_paragraph().paragraph_format.space_after = Pt(2)

            elif b.kind == "list":
                for depth, item in b.rows:
                    style = "List Number" if b.ordered else "List Bullet"
                    try:
                        p = doc.add_paragraph(style=style)
                    except KeyError:
                        p = doc.add_paragraph()
                    p.paragraph_format.left_indent = Inches(0.25 + 0.22 * depth)
                    p.paragraph_format.space_after = Pt(2)
                    add_runs(p, item, base_size=10)

            elif b.kind == "hr":
                p = doc.add_paragraph()
                r = p.add_run("─" * 58)
                r.font.color.rgb = RGBColor.from_string(RULE)
                r.font.size = Pt(8)

    doc.save(path)


# ─────────────────────────────────────────────────────── 3 · pdf ──

def build_pdf(path: str) -> None:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm, mm
    from reportlab.lib import colors
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.enums import TA_CENTER, TA_LEFT
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.platypus import (
        BaseDocTemplate, PageTemplate, Frame, Paragraph, Spacer, Table,
        TableStyle, PageBreak, KeepTogether, Flowable,
    )
    from reportlab.platypus.tableofcontents import TableOfContents

    fonts = r"C:\Windows\Fonts"

    def reg(name, filename, fallback):
        try:
            pdfmetrics.registerFont(TTFont(name, os.path.join(fonts, filename)))
            return name
        except Exception:
            return fallback

    BODY = reg("Body", "calibri.ttf", "Helvetica")
    BODY_B = reg("BodyB", "calibrib.ttf", "Helvetica-Bold")
    BODY_I = reg("BodyI", "calibrii.ttf", "Helvetica-Oblique")
    MONO = reg("Mono", "consola.ttf", "Courier")
    MONO_B = reg("MonoB", "consolab.ttf", "Courier-Bold")

    try:
        from reportlab.pdfbase.pdfmetrics import registerFontFamily
        registerFontFamily(BODY, normal=BODY, bold=BODY_B, italic=BODY_I, boldItalic=BODY_B)
        registerFontFamily(MONO, normal=MONO, bold=MONO_B, italic=MONO, boldItalic=MONO_B)
    except Exception:
        pass

    C_INK = colors.HexColor("#" + INK)
    C_MUTED = colors.HexColor("#" + MUTED)
    C_ACCENT = colors.HexColor("#" + ACCENT)
    C_ACCENT_DK = colors.HexColor("#" + ACCENT_DK)
    C_RULE = colors.HexColor("#" + RULE)
    C_CODE = colors.HexColor("#" + CODE_BG)
    C_BAND = colors.HexColor("#" + BAND)

    PW, PH = A4
    LM = RM = 1.5 * cm
    TM = 1.7 * cm
    BM = 1.6 * cm
    AVAIL = PW - LM - RM

    ss = getSampleStyleSheet()

    def st(name, **kw):
        base = dict(fontName=BODY, fontSize=9.5, leading=13.2, textColor=C_INK,
                    spaceBefore=0, spaceAfter=5, alignment=TA_LEFT)
        base.update(kw)
        return ParagraphStyle(name, **base)

    S_H1 = st("H1", fontName=BODY_B, fontSize=19, leading=23, textColor=C_ACCENT_DK,
              spaceBefore=6, spaceAfter=10)
    S_H2 = st("H2", fontName=BODY_B, fontSize=14, leading=18, textColor=C_ACCENT_DK,
              spaceBefore=13, spaceAfter=6)
    S_H3 = st("H3", fontName=BODY_B, fontSize=11.5, leading=15, textColor=C_INK,
              spaceBefore=10, spaceAfter=4)
    S_H4 = st("H4", fontName=BODY_B, fontSize=10, leading=13.5, textColor=C_MUTED,
              spaceBefore=8, spaceAfter=3)
    S_BODY = st("Body")
    S_QUOTE = st("Quote", fontName=BODY_I, fontSize=9.2, leading=13,
                 textColor=C_ACCENT_DK, leftIndent=9, spaceBefore=3, spaceAfter=6)
    S_LIST = st("List", leftIndent=12, spaceAfter=2.5)
    S_CELL = st("Cell", fontSize=8.2, leading=10.6, spaceAfter=0)
    S_CELLH = st("CellH", fontName=BODY_B, fontSize=8.2, leading=10.6,
                 spaceAfter=0, textColor=colors.white)
    S_CAP = st("Cap", fontName=BODY_B, fontSize=7.6, leading=9.6,
               textColor=C_ACCENT, spaceAfter=1.5)
    S_COVER_T = st("CoverT", fontName=BODY_B, fontSize=33, leading=38,
                   textColor=C_ACCENT_DK, alignment=TA_CENTER, spaceAfter=6)
    S_COVER_S = st("CoverS", fontSize=16, leading=21, textColor=C_MUTED,
                   alignment=TA_CENTER, spaceAfter=16)
    S_COVER_M = st("CoverM", fontSize=10.5, leading=15, textColor=C_MUTED,
                   alignment=TA_CENTER, spaceAfter=3)
    S_TOC1 = st("TOC1", fontName=BODY_B, fontSize=11, leading=18, spaceBefore=7)
    S_TOC2 = st("TOC2", fontSize=9.4, leading=14.5, leftIndent=16)
    S_TOC3 = st("TOC3", fontSize=8.6, leading=13, leftIndent=32, textColor=C_MUTED)

    def esc(text: str) -> str:
        return (text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"))

    def rl_inline(text: str) -> str:
        text = strip_glyphs(text)
        out = []
        for chunk, bold, italic, code, href in inline_tokens(text):
            if not chunk:
                continue
            piece = esc(chunk)
            if code:
                piece = f'<font face="{MONO}" size="8.4" color="#{BAD}">{piece}</font>'
            if bold:
                piece = f"<b>{piece}</b>"
            if italic:
                piece = f"<i>{piece}</i>"
            if href:
                safe = href.replace("&", "&amp;").replace('"', "%22")
                piece = f'<link href="{safe}" color="#{ACCENT}">{piece}</link>'
            out.append(piece)
        return "".join(out) or "&nbsp;"

    # ── vector diagrams drawn natively ──
    class Diagram(Flowable):
        """Native vector drawing — sharp at any zoom, no raster assets."""

        def __init__(self, kind, width, height):
            super().__init__()
            self.kind = kind
            self.width = width
            self.height = height

        def wrap(self, aw, ah):
            return (self.width, self.height)

        def _box(self, x, y, w, h, label, sub=None, fill="#FFFFFF",
                 stroke="#1D4ED8", tcol="#1F2933", radius=5, bold=True):
            c = self.canv
            c.saveState()
            c.setFillColor(colors.HexColor(fill))
            c.setStrokeColor(colors.HexColor(stroke))
            c.setLineWidth(1.1)
            c.roundRect(x, y, w, h, radius, stroke=1, fill=1)
            c.setFillColor(colors.HexColor(tcol))
            c.setFont(BODY_B if bold else BODY, 8.4)
            if sub:
                c.drawCentredString(x + w / 2, y + h / 2 + 3.4, label)
                c.setFont(BODY, 7)
                c.setFillColor(colors.HexColor("#" + MUTED))
                c.drawCentredString(x + w / 2, y + h / 2 - 6.4, sub)
            else:
                c.drawCentredString(x + w / 2, y + h / 2 - 3, label)
            c.restoreState()

        def _arrow(self, x1, y1, x2, y2, label=None, col="#5A6875", dashed=False,
                   head=True):
            import math
            c = self.canv
            c.saveState()
            c.setStrokeColor(colors.HexColor(col))
            c.setLineWidth(1.0)
            if dashed:
                c.setDash(2.5, 2.5)
            c.line(x1, y1, x2, y2)
            c.setDash()
            if head:
                ang = math.atan2(y2 - y1, x2 - x1)
                spread = 0.40           # ~23 degrees each side
                size = 5.2
                c.setFillColor(colors.HexColor(col))
                p = c.beginPath()
                p.moveTo(x2, y2)
                p.lineTo(x2 - size * math.cos(ang - spread),
                         y2 - size * math.sin(ang - spread))
                p.lineTo(x2 - size * math.cos(ang + spread),
                         y2 - size * math.sin(ang + spread))
                p.close()
                c.drawPath(p, stroke=0, fill=1)
            if label:
                c.setFillColor(colors.HexColor(col))
                c.setFont(BODY, 6.6)
                c.drawCentredString((x1 + x2) / 2, (y1 + y2) / 2 + 3.6, label)
            c.restoreState()

        def draw(self):
            getattr(self, "_draw_" + self.kind)()

        def _draw_context(self):
            W, H = self.width, self.height
            cy = H - 112
            bw, bh = W * 0.56, 34
            x_app = (W - bw) / 2

            up = [("17 SKUs of hardware", "ESP32 \u00b7 MQTT/TLS \u00b7 own CA",
                   "#FEF3C7", "#B45309"),
                  ("platform/ control plane", "ONE Oracle free-tier VM",
                   "#F5F3FF", "#6D28D9")]
            uw = (W - 90) / 2
            up_y = H - 34
            u_centres = []
            for n, (lab, sub, fill, stroke) in enumerate(up):
                x = n * (uw + 90)
                self._box(x, up_y, uw, 28, lab, sub, fill=fill, stroke=stroke)
                u_centres.append(x + uw / 2)

            bus_y = up_y - 14
            c = self.canv
            c.saveState()
            c.setStrokeColor(colors.HexColor("#64748B"))
            c.setLineWidth(0.9)
            c.line(u_centres[0], bus_y, u_centres[-1], bus_y)
            for cx in u_centres:
                c.line(cx, up_y, cx, bus_y)
            c.restoreState()
            self._arrow(W / 2, bus_y, W / 2, cy + bh + 4,
                        "telemetry   \u00b7   federated sign-in", col="#1D4ED8")

            self._box(x_app, cy, bw, bh, "website  (circuvent-technologies)",
                      "5 products \u00b7 150 API routes \u00b7 428,355 lines",
                      fill="#DBEAFE", stroke="#1D4ED8")

            row_y = cy - 84
            deps = [("Neon Postgres", "10 tables \u00b7 made at boot",
                     "#EFF6FF", "#1D4ED8"),
                    (".data/ JSON files", "~27 modules \u00b7 memory-only",
                     "#FEF2F2", "#B91C1C"),
                    ("Razorpay \u00b7 SMTP", "webhook is a STUB",
                     "#ECFDF5", "#15803D")]
            gap = 8.0
            dw = (W - gap * (len(deps) - 1)) / len(deps)
            d_centres = []
            for n, (lab, sub, fill, stroke) in enumerate(deps):
                x = n * (dw + gap)
                self._box(x, row_y, dw, 34, lab, sub, fill=fill, stroke=stroke)
                d_centres.append((x + dw / 2, stroke))

            dbus = cy - 22
            c = self.canv
            c.saveState()
            c.setStrokeColor(colors.HexColor("#1D4ED8"))
            c.setLineWidth(1.0)
            c.line(W / 2, cy - 3, W / 2, dbus)
            c.line(d_centres[0][0], dbus, d_centres[-1][0], dbus)
            c.restoreState()
            for cx, col in d_centres:
                self._arrow(cx, dbus, cx, row_y + 34 + 4, col=col)

            c = self.canv
            c.saveState()
            c.setFont(BODY, 6.6)
            c.setFillColor(colors.HexColor("#" + MUTED))
            c.drawCentredString(W / 2, row_y - 13,
                                "It is called \"website\". It is a company: a storefront, an IoT "
                                "platform, firmware for a 17-product line, PCB designs and a "
                                "shipped mobile app.")
            c.restoreState()

        def _draw_layers(self):
            W, H = self.width, self.height
            rows = [
                ("THE GATE", "src/proxy.ts - host mounts, redirects, CSP - AND NO AUTHENTICATION",
                 "#FEF2F2", "#B91C1C"),
                ("ROUTES", "150 handlers - each enforces its own auth, five schemes coexist",
                 "#EFF6FF", "#1D4ED8"),
                ("AUTHORIZATION", "AdminRole x AdminArea capability gate - 5 roles, ~20 areas",
                 "#ECFDF5", "#15803D"),
                ("DOMAIN MODULES", "288 lib files - icm, commerce, smarthome, telemetry, reports",
                 "#FFF7ED", "#B45309"),
                ("STORAGE CHOICE", "db.ts (10 tables) OR data-file.ts (JSON) - only 3 of ~30 durable",
                 "#F5F3FF", "#6D28D9"),
                ("POSTGRES", "Neon HTTP driver - no transactions, no RLS, no foreign keys",
                 "#F1F5F9", "#5A6875"),
            ]
            h = 26
            gap = 6
            y = H - h
            for lab, sub, fill, stroke in rows:
                self._box(0, y, W, h, "", None, fill=fill, stroke=stroke)
                c = self.canv
                c.saveState()
                c.setFillColor(colors.HexColor(stroke))
                c.setFont(BODY_B, 8.0)
                c.drawString(9, y + h / 2 + 2.4, lab)
                c.setFillColor(colors.HexColor("#" + MUTED))
                c.setFont(BODY, 6.9)
                c.drawString(9, y + h / 2 - 7.4, sub)
                c.restoreState()
                y -= (h + gap)

            c = self.canv
            c.saveState()
            c.setFont(BODY, 6.6)
            c.setFillColor(colors.HexColor("#" + MUTED))
            c.drawCentredString(W / 2, y + h - 4,
                                "Fourteen CI steps would check this stack. "
                                "Twenty-seven runs have failed to start.")
            c.restoreState()

        def _draw_apply(self):
            W, H = self.width, self.height
            lanes = ["Customer", "website", "Razorpay", "Datastore", "Email"]
            n = len(lanes)
            lw = W / n
            c = self.canv
            top = H - 16
            for i, lab in enumerate(lanes):
                x = i * lw + lw / 2
                c.saveState()
                c.setFillColor(colors.HexColor("#" + ACCENT_DK))
                c.setFont(BODY_B, 7.4)
                c.drawCentredString(x, top, lab)
                c.setStrokeColor(colors.HexColor("#" + RULE))
                c.setLineWidth(0.8)
                c.setDash(2, 3)
                c.line(x, top - 8, x, 10)
                c.setDash()
                c.restoreState()
            steps = [
                (0, 1, "checkout"),
                (1, 1, "priceItems() recomputes from the catalogue"),
                (1, 1, "client-submitted prices are NEVER trusted"),
                (1, 2, "Math.round(due * 100) - paise - create order"),
                (2, 0, "payment UI"),
                (0, 1, "handback signature"),
                (1, 1, "verify HMAC"),
                (1, 2, "RE-FETCH the payment from the gateway"),
                (2, 1, "require status === 'captured' before trusting it"),
                (1, 3, "mark paid"),
                (1, 4, "receipt - awaited inline, not queued"),
                (2, 1, "webhook: signature verified, then NOTHING HAPPENS"),
            ]
            y = top - 22
            for a, b, lab in steps:
                xa = a * lw + lw / 2
                xb = b * lw + lw / 2
                if a == b:
                    c.saveState()
                    c.setStrokeColor(colors.HexColor("#" + WARN))
                    c.setLineWidth(1.0)
                    c.rect(xa, y - 4, 30, 10, stroke=1, fill=0)
                    c.setFillColor(colors.HexColor("#" + WARN))
                    c.setFont(BODY, 6.3)
                    c.drawString(xa + 34, y - 1, lab)
                    c.restoreState()
                else:
                    self._arrow(xa, y, xb, y, lab,
                                col="#1D4ED8" if b > a else "#15803D")
                y -= 15.6

    # ── flowable builders ──
    def code_flowable(block):
        # Measure the real rendered width instead of guessing from the
        # character count: box-drawing glyphs get substituted before
        # layout, and a hand-tuned per-character estimate under-measures
        # and silently clipped wide ASCII diagrams at the right margin.
        from reportlab.pdfbase.pdfmetrics import stringWidth
        measured = [strip_glyphs(x, mono=True) for x in block.lines] or [""]
        size = 7.6
        while size > 4.2:
            widest = max((stringWidth(x, MONO, size) for x in measured),
                         default=0.0)
            if widest <= AVAIL - 14:
                break
            size -= 0.2
        leading = size * 1.32
        is_mermaid = block.lang.lower() == "mermaid"

        frame_h = PH - TM - BM
        # Leave room for the caption, padding and a little slack, then work out
        # how many lines of this size fit on one page. A single un-splittable
        # table taller than the frame is a hard LayoutError, so long blocks —
        # the big Mermaid class diagram in particular — are chunked instead.
        per_page = max(8, int((frame_h - 60) / leading))

        chunks = [block.lines[i:i + per_page]
                  for i in range(0, len(block.lines), per_page)] or [[""]]

        parts = []
        for n, chunk in enumerate(chunks):
            body = "\n".join(strip_glyphs(x, mono=True) or " " for x in chunk)
            para = Paragraph(
                esc(body).replace("\n", "<br/>").replace(" ", "&nbsp;"),
                ParagraphStyle("Code", fontName=MONO, fontSize=size,
                               leading=leading, textColor=C_INK,
                               spaceBefore=0, spaceAfter=0),
            )
            t = Table([[para]], colWidths=[AVAIL])
            t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, -1),
                 colors.HexColor("#EDF2F7") if is_mermaid else C_CODE),
                ("BOX", (0, 0), (-1, -1), 0.6,
                 C_ACCENT if is_mermaid else C_RULE),
                ("LEFTPADDING", (0, 0), (-1, -1), 6),
                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                ("TOPPADDING", (0, 0), (-1, -1), 5),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ]))
            if n == 0:
                if is_mermaid:
                    parts.append(Paragraph(
                        "Mermaid diagram &mdash; renders interactively in the "
                        "Markdown edition", S_CAP))
                parts.append(t)
            else:
                parts.append(PageBreak())
                parts.append(Paragraph(
                    f"(continued &mdash; part {n + 1} of {len(chunks)})", S_CAP))
                parts.append(t)
        parts.append(Spacer(1, 6))
        return parts

    def table_flowable(block):
        header, *rows = block.rows
        ncols = len(header)
        # A Markdown table may legitimately have an empty header row — several
        # in these documents are two-column fact lists. Rendering the dark
        # header band for one produces a blank blue bar, so drop it instead.
        has_header = any(c.strip() for c in header)
        if not has_header:
            header = []
        weights = []
        for i in range(ncols):
            col = ([header[i]] if has_header else []) + [r[i] for r in rows]
            weights.append(max(6, min(60, max((len(str(x)) for x in col),
                                              default=6))))
        total = float(sum(weights)) or 1.0
        widths = [max(38.0, AVAIL * w / total) for w in weights]
        scale = AVAIL / sum(widths)
        widths = [w * scale for w in widths]

        data = []
        if has_header:
            data.append([Paragraph(rl_inline(c), S_CELLH) for c in header])
        for r in rows:
            data.append([Paragraph(rl_inline(c), S_CELL) for c in r])
        if not data:
            return []

        t = Table(data, colWidths=widths,
                  repeatRows=1 if has_header else 0, hAlign="LEFT")
        style = [
            ("GRID", (0, 0), (-1, -1), 0.4, C_RULE),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 4.5),
            ("RIGHTPADDING", (0, 0), (-1, -1), 4.5),
            ("TOPPADDING", (0, 0), (-1, -1), 3.5),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 3.5),
        ]
        if has_header:
            style += [("BACKGROUND", (0, 0), (-1, 0), C_ACCENT_DK),
                      ("TEXTCOLOR", (0, 0), (-1, 0), colors.white)]
        else:
            style.append(("BACKGROUND", (0, 0), (0, -1),
                          colors.HexColor("#F1F5F9")))
        start = 1 if has_header else 0
        for i in range(start, len(data)):
            if (i - start) % 2 == 1:
                style.append(("BACKGROUND", (0, i), (-1, i),
                              colors.HexColor("#F8FAFC")))
        t.setStyle(TableStyle(style))
        return [t, Spacer(1, 7)]

    def quote_flowable(lines):
        inner = [Paragraph(rl_inline(x), S_QUOTE) for x in lines]
        t = Table([[inner]], colWidths=[AVAIL])
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#EEF3FB")),
            ("LINEBEFORE", (0, 0), (0, -1), 2.4, C_ACCENT),
            ("LEFTPADDING", (0, 0), (-1, -1), 9),
            ("RIGHTPADDING", (0, 0), (-1, -1), 8),
            ("TOPPADDING", (0, 0), (-1, -1), 5),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ]))
        return [t, Spacer(1, 6)]

    # ── document scaffolding ──
    class Doc(BaseDocTemplate):
        def afterFlowable(self, flowable):
            if not isinstance(flowable, Paragraph):
                return
            style = flowable.style.name
            level = {"H1": 0, "H2": 1, "H3": 2}.get(style)
            if level is None:
                return
            text = re.sub(r"<[^>]+>", "", flowable.getPlainText())
            self.notify("TOCEntry", (level, text, self.page))

    def decorate(canv, doc_):
        canv.saveState()
        if doc_.page > 1:
            canv.setFillColor(C_BAND)
            canv.rect(0, PH - 13 * mm, PW, 13 * mm, stroke=0, fill=1)
            canv.setFillColor(colors.white)
            canv.setFont(BODY_B, 7.6)
            canv.drawString(LM, PH - 8.7 * mm, TITLE)
            canv.setFont(BODY, 7.6)
            canv.drawRightString(PW - RM, PH - 8.7 * mm, SUBTITLE)

            canv.setStrokeColor(C_RULE)
            canv.setLineWidth(0.6)
            canv.line(LM, BM - 4, PW - RM, BM - 4)
            canv.setFillColor(C_MUTED)
            canv.setFont(BODY, 7.2)
            canv.drawString(LM, BM - 13, f"{ORG} · generated {BUILT}")
            canv.drawRightString(PW - RM, BM - 13, f"Page {doc_.page}")
        canv.restoreState()

    doc = Doc(path, pagesize=A4, leftMargin=LM, rightMargin=RM,
              topMargin=TM, bottomMargin=BM,
              title=f"{TITLE} — {SUBTITLE}", author=ORG,
              subject="Architecture and technical audit")
    frame = Frame(LM, BM, AVAIL, PH - TM - BM, id="body",
                  leftPadding=0, rightPadding=0, topPadding=0, bottomPadding=0)
    doc.addPageTemplates([PageTemplate(id="main", frames=[frame], onPage=decorate)])

    story = []

    # cover
    story.append(Spacer(1, 3.6 * cm))
    story.append(Paragraph(TITLE, S_COVER_T))
    story.append(Paragraph(SUBTITLE, S_COVER_S))
    rule = Table([[""]], colWidths=[6 * cm], rowHeights=[2.2])
    rule.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, -1), C_ACCENT)]))
    rule.hAlign = "CENTER"
    story.append(rule)
    story.append(Spacer(1, 0.8 * cm))
    story.append(Paragraph(ORG, S_COVER_M))
    story.append(Paragraph("Next.js 16 · Neon · ESP32 · MQTT/TLS · Expo · KiCad",
                           S_COVER_M))
    story.append(Paragraph(f"Generated {BUILT}", S_COVER_M))
    story.append(Spacer(1, 1.1 * cm))
    story.append(Diagram("context", AVAIL, 232))
    story.append(Spacer(1, 0.7 * cm))

    facts = [
        ["Architecture", "Five web products, firmware, hardware and two native apps in one repository"],
        ["Scale", "23,184 files · 1,810 src TypeScript · 428,355 lines · 150 API routes · 108 pages"],
        ["Database", "Neon Postgres, HTTP driver · 10 tables · CREATED AT RUNTIME · no migrations"],
        ["Sessions", "Home-grown HMAC-SHA256 — no JWT library is used anywhere in src/"],
        ["Devices", "17 retail SKUs · ESP32 · MQTT/TLS on an own broker with an own CA"],
        ["Tests", "236 suites · 4,328 tests · 1 failing · 49 seconds"],
        ["CI/CD", "14 steps, 13 hard gates — and 27 of 27 runs failed to start"],
        ["Primary risk", "~27 storage modules are memory-only in production, passkeys among them"],
    ]
    ft = Table([[Paragraph(f"<b>{a}</b>", S_CELL), Paragraph(b, S_CELL)]
                for a, b in facts], colWidths=[AVAIL * 0.32, AVAIL * 0.68])
    ft.setStyle(TableStyle([
        ("GRID", (0, 0), (-1, -1), 0.4, C_RULE),
        ("BACKGROUND", (0, 0), (0, -1), colors.HexColor("#F1F5F9")),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
    ]))
    story.append(ft)
    story.append(PageBreak())

    # contents
    story.append(Paragraph("Contents", S_H1))
    toc = TableOfContents()
    toc.levelStyles = [S_TOC1, S_TOC2, S_TOC3]
    story.append(toc)
    story.append(PageBreak())

    # body
    for idx, (fname, label) in enumerate(SOURCES, start=1):
        with open(os.path.join(HERE, fname), "r", encoding="utf-8") as fh:
            blocks = parse_markdown(fh.read())

        if idx > 1:
            story.append(PageBreak())

        band = Table([[Paragraph(
            f'<font color="white" size="8"><b>PART {idx}</b></font>', S_CELL)]],
            colWidths=[AVAIL])
        band.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, -1), C_BAND),
            ("LEFTPADDING", (0, 0), (-1, -1), 8),
            ("TOPPADDING", (0, 0), (-1, -1), 3),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
        ]))
        story.append(band)
        story.append(Spacer(1, 8))

        if idx == 1:
            pass

        for b in blocks:
            if b.kind == "heading":
                style = {1: S_H1, 2: S_H2, 3: S_H3}.get(b.level, S_H4)
                story.append(Paragraph(rl_inline(b.text), style))
                # insert the native diagrams alongside the sections they explain
                if idx == 1 and b.level == 2 and b.text.startswith("4."):
                    story.append(Spacer(1, 3))
                    story.append(Paragraph("Figure — module and layer map", S_CAP))
                    story.append(Diagram("layers", AVAIL, 200))
                    story.append(Spacer(1, 8))
                if idx == 1 and b.level == 2 and b.text.startswith("7."):
                    story.append(Spacer(1, 3))
                    story.append(Paragraph("Figure — checkout and payment, "
                                           "end to end", S_CAP))
                    story.append(Diagram("apply", AVAIL, 214))
                    story.append(Spacer(1, 8))

            elif b.kind == "para":
                story.append(Paragraph(rl_inline(b.text), S_BODY))

            elif b.kind == "quote":
                story.extend(quote_flowable(b.lines))

            elif b.kind == "code":
                story.extend(code_flowable(b))

            elif b.kind == "table":
                story.extend(table_flowable(b))

            elif b.kind == "list":
                counter = 0
                for depth, item in b.rows:
                    if b.ordered:
                        counter += 1
                        marker = f"<b>{counter}.</b>"
                    else:
                        marker = "&bull;"
                    story.append(Paragraph(
                        f"{marker}&nbsp;&nbsp;{rl_inline(item)}",
                        ParagraphStyle(f"L{depth}", parent=S_LIST,
                                       leftIndent=12 + 14 * depth)))
                story.append(Spacer(1, 4))

            elif b.kind == "hr":
                story.append(Spacer(1, 3))
                hr = Table([[""]], colWidths=[AVAIL], rowHeights=[0.7])
                hr.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, -1), C_RULE)]))
                story.append(hr)
                story.append(Spacer(1, 6))

    doc.multiBuild(story)


# ────────────────────────────────────────────────────── 4 · pptx ──

SLIDES = [
    {
        "kind": "bullets",
        "title": "What this repository is",
        "lead": "It is called \"website\". It is a company.",
        "bullets": [
            ("A marketing site, an e-commerce store with payments, an "
             "83-route admin back office, a smart-home IoT console with 60+ "
             "sections, and a public developer portal — all in one Next.js "
             "application.", 0),
            ("Plus firmware for a 17-SKU retail hardware line sold on "
             "Amazon.in and Flipkart, real KiCad PCB designs, a self-hosted "
             "MQTT cloud on one Oracle VM, a shipped Play Store app, a "
             "drone, and an RC car.", 0),
            ("Plus a completely separate internal SaaS for project "
             "tracking, HR, payroll and a client portal.", 0),
            ("23,184 files · 1,810 src TypeScript files · 428,355 lines · "
             "150 API routes · 108 pages · 533 commits · 5 git remotes.", 0),
            ("The README describes 18 routes and 50+ components. That was "
             "true once.", 0),
        ],
    },
    {
        "kind": "facts",
        "title": "At a glance",
        "facts": [
            ("Stack", "Next.js 16.2 App Router · React 19.2 · TypeScript 5 strict"),
            ("Database", "Neon Postgres, HTTP driver · 10 tables · created at RUNTIME"),
            ("Sessions", "Home-grown HMAC-SHA256. No JWT library is used at all."),
            ("Devices", "17 SKUs · ESP32 · MQTT/TLS on an own broker with an own CA"),
            ("Mobile", "Expo, SHIPPED at v1.13.1 · a Kotlin/Swift prototype alongside"),
            ("Tests", "236 suites · 4,328 tests · 1 failing · 49 seconds"),
            ("CI/CD", "14 steps, 13 hard gates — and 27 of 27 runs failed to start"),
            ("Observability", "None. No error tracking, no alerting, no log sink."),
        ],
    },
    {
        "kind": "topology",
        'top': [('17 SKUs of hardware', 'ESP32 · MQTT/TLS', 'warn'), ('website', '150 routes · 5 products', 'primary'), ('platform/', 'ONE Oracle VM', 'violet')],
        'links': ['telemetry', 'federation'],
        'hub': 1,
        'children': [('Neon Postgres', '10 tables · made at boot'), ('.data/ JSON files', '~27 modules, memory-only'), ('Razorpay', 'webhook is a STUB'), ('mobile/ Expo', 'SHIPPED v1.13.1')],
        'childnote': 'Wi-Fi, GSM, LoRa, ESP-NOW and MAVLink all terminate on one free-tier virtual machine.',
        "title": "Architectural topology",
        "note": "circuvent-platform/ is a DIFFERENT product entirely. The names are a trap.",
    },
    {
        "kind": "bullets",
        "title": "The finding that matters most",
        "lead": "Roughly twenty-seven storage modules have no database behind them.",
        "bullets": [
            ("createFileStore() keeps an in-memory working copy and writes "
             "through to a JSON file. Its own comment: \u201cOn read-only "
             "filesystems... writes silently stop and the module degrades to "
             "in-memory-only for that instance, instead of throwing and "
             "breaking the request.\u201d", 0),
            ("About thirty modules are built on it. EXACTLY THREE pass "
             "durable: true — icm-store, admin-warranty and api-failures.", 0),
            ("The other ~27 include CMS content, CRM records, pricing, "
             "currency, tax configuration, feature flags, marketing, staff "
             "activity, 570 KB of telemetry, developer-portal tokens — and "
             "PASSKEYS.", 0),
            ("In production those live in one lambda instance's memory. When "
             "the instance recycles the data is gone. Not stale. Not "
             "corrupted. Gone.", 0),
            ("AND THE REPOSITORY ALREADY KNOWS: \u201cIncidents filed weeks "
             "ago were not hidden; they were gone.\u201d That bug was found, "
             "understood, written down — and fixed for ONE module.", 0),
        ],
    },
    {
        "kind": "layers",
        "title": "Layered design",
        "note": "Auth is enforced per route, by helper calls. The gate does not gate.",
    },
    {
        "kind": "flow",
        "title": "Data flow - schema, and where it comes from",
        "steps": [
            ("1", "Cold start", "A lambda boots. Nothing exists yet."),
            ("2", "initDb()", "Runs an array of CREATE TABLE IF NOT EXISTS strings"),
            ("3", "Every call", "Every db* function awaits initDb() before querying"),
            ("4", "Consequence", "The role must hold DDL rights, permanently"),
            ("5", "Driver", "neon() HTTP — so transactions are IMPOSSIBLE"),
            ("6", "store_kv", "23 shop collections, ONE JSONB ROW EACH, key='_all'"),
            ("7", "Fallback", "No DATABASE_URL? Silently write JSON files instead"),
            ("8", "Production", "Read-only disk. The write fails. The catch swallows it."),
        ],
    },
    {
        "kind": "bullets",
        "title": "Firmware: no signature verification",
        "lead": "On devices that switch mains relays and door locks.",
        "bullets": [
            ("Both OTA paths download over a certificate-PINNED TLS "
             "connection. That authenticates the SERVER. It does not "
             "authenticate the FIRMWARE.", 0),
            ("Searched for and not found anywhere: Ed25519 or RSA signature "
             "verification, a hash manifest checked against an on-device "
             "public key, or ESP32 Secure Boot in any of the 29 "
             "platformio.ini files.", 0),
            ("The code articulates the stakes itself, explaining why "
             "setInsecure() was removed: \u201canyone able to intercept that "
             "connection... could serve arbitrary firmware and take permanent "
             "control of a board that switches mains relays and door "
             "locks.\u201d", 0),
            ("The transport hole was closed. The integrity hole was not.", 0),
            ("And the device polls a manifest endpoint on a timer that DOES "
             "NOT EXIST — platform/api has no firmware route. The hardware "
             "checklist marks both signed builds and a rollback plan "
             "unchecked.", 0),
        ],
    },
    {
        "kind": "bullets",
        "title": "The payment webhook is a stub",
        "lead": "It verifies its signature correctly. Then it does nothing.",
        "bullets": [
            ("const expected = crypto.createHmac(\"sha256\", secret)"
             ".update(raw).digest(\"hex\"); — correct, and constant-time "
             "compared.", 0),
            ("Then it console.log()s and returns. Its own comment: "
             "\u201cReconciliation hook: when a persistent order store "
             "exists, mark the order paid/failed here...\u201d", 0),
            ("Any payment captured at the gateway whose browser never "
             "returns — a closed tab, a dropped connection, a failed "
             "redirect — is money taken and an order never marked paid.", 0),
            ("The checkout-side path is genuinely good: verifyCapturedPayment "
             "re-fetches the payment from Razorpay and requires status === "
             "\u201ccaptured\u201d before trusting the amount, so a forged "
             "client signature cannot credit an order.", 0),
            ("The webhook was meant to be the safety net under that.", 0),
        ],
    },
    {
        "kind": "bullets",
        "title": "Five credential schemes, one gate that doesn't gate",
        "lead": "src/proxy.ts performs no authentication at all.",
        "bullets": [
            ("Customer sessions: HMAC-SHA256 over email|issuedAt|"
             "tokenVersion, signed with ACCOUNT_SECRET.", 0),
            ("Staff sessions: the same shape, but ADMIN_SECRET — \u201cStaff "
             "sessions get their own key so a leaked customer key cannot mint "
             "one.\u201d", 0),
            ("Plus a static ADMIN_API_KEY header, a CRON_SECRET bearer, and a "
             "bridged external JWT for the control plane. Zero use of jose or "
             "jsonwebtoken anywhere in src/.", 0),
            ("This app does NOT share the suite's HS256 AUTH_JWT_SECRET. It "
             "is the only Circuvent application outside that scheme.", 0),
            ("proxy.ts handles host mounts, a legacy redirect, X-Request-Id "
             "and the CSP header — and no auth. Seven routes could not be "
             "matched to any known mechanism at all.", 0),
        ],
    },
    {
        "kind": "bullets",
        "title": "The same rule, written four times",
        "lead": "And the bug it causes has already shipped twice.",
        "bullets": [
            ("The per-device capability table decides which fields a device "
             "reports, which command keys control it, and which device types "
             "must expose no toggle at all — a camera's only switch is "
             "streaming; a drone's is flight permission.", 0),
            ("It is implemented independently in Expo TypeScript, in Kotlin, "
             "in Swift, and again in the web console. No shared schema, no "
             "code generation.", 0),
            ("native/README.md names the failure: \u201cA Home Hub reports "
             "power2 and is commanded with {ch: 1, on: true}... Send the "
             "state key to a device that wanted the command key and nothing "
             "errors. THAT BUG HAS ALREADY SHIPPED TWICE — once on the web "
             "and once in the Expo app.\u201d", 0),
            ("The only guard is a parity test that diffs the three tables. "
             "Drift is CAUGHT BY AN ASSERTION, not PREVENTED BY A SCHEMA.", 0),
        ],
    },
    {
        "kind": "bullets",
        "title": "CI has never executed",
        "lead": "The most thorough pipeline in the suite. Zero runs.",
        "bullets": [
            ("27 workflow runs on record. 27 startup_failure. Zero seconds "
             "each. And the registered \u201cCI\u201d workflow has zero runs "
             "attributed to it at all — every one belongs to a synthetic, "
             "deleted placeholder.", 0),
            ("What is not running: a typecheck, 4,328 root tests, the control "
             "plane's own ~290 tests, a PGlite database test, a production "
             "build, and a full Playwright suite. Thirteen hard gates; only "
             "lint is advisory.", 0),
            ("The workflow's own comments describe fixing exactly the gaps "
             "that made those tests unable to block a deploy: \u201cThat left "
             "~290 tests... unable to block a deploy\u201d and \u201cE2E was "
             "never run in CI, which is how the sitemap assertion stayed "
             "broken.\u201d Those fixes have never run either.", 0),
            ("And verify:secrets is not a CI gate — it is guarded only by a "
             "local pre-commit hook that git commit --no-verify defeats.", 0),
        ],
    },
    {
        "kind": "bullets",
        "title": "Every module names the bug it prevents",
        "lead": "The best convention in the entire Circuvent suite.",
        "bullets": [
            ("passkeys.ts: \u201cThe passkey still existed, still verified, "
             "and belonged to nobody.\u201d", 0),
            ("sso.ts: \u201cProduction users could sign in to dev, and dev "
             "quietly accumulated live credentials while doing it. The "
             "isolation guard was pointed at the wrong door.\u201d", 0),
            ("install-hooks.mjs: \u201cbash on the Windows box that does the "
             "builds is WSL with no distribution installed — it failed "
             "silently and the build carried on with the wrong signing "
             "key.\u201d", 0),
            ("check-no-secrets.js: \u201ca secret pushed once is in every "
             "clone and every fork, and deleting the file in a later commit "
             "does not remove it from history.\u201d", 0),
            ("verify_kt_docs.py: \u201ca build script in this repository has "
             "previously reported success while publishing the previous run's "
             "artifact, so 'the command succeeded' is not evidence.\u201d", 0),
        ],
    },
    {
        "kind": "bullets",
        "title": "Verification that proves, not asserts",
        "lead": "The engineering instinct here is excellent. Almost none of it is automated.",
        "bullets": [
            ("verify-icm-durability.ts spawns FOUR REAL OPERATING-SYSTEM "
             "PROCESSES, sharing no memory, to prove data survives a cold "
             "start — rather than mocking one. It was run during this audit "
             "and passed.", 0),
            ("audit-code-contrast.mjs reads COMPUTED CSS in a real browser, "
             "\u201cbecause the bug it checks for was invisible to every "
             "existing test... the only symptom was that a human could not "
             "read it.\u201d", 0),
            ("perf-probe.mjs uses the Resource Timing API for real "
             "transferred bytes and DIRECTLY OBSERVES layout shift rather "
             "than inferring it from missing dimensions.", 0),
            ("verify_business_docs.py opens the generated PowerPoint and "
             "asserts a real catalogue price is inside it — and that no raw "
             "unformatted price slipped past the formatter.", 0),
            ("None of these five scripts runs in CI. Thirteen npm scripts "
             "run nowhere automated at all.", 0),
        ],
    },
    {
        "kind": "facts",
        "title": "The sub-projects, sized honestly",
        "facts": [
            ("firmware/", "13,003 files on disk — but git tracks only 84"),
            ("", "29 sketches, one shared library, 17 retail products"),
            ("platform/", "The IoT control plane. Mosquitto + Express + Postgres."),
            ("circuvent-platform/", "A DIFFERENT product: internal HR, payroll, projects"),
            ("mobile/", "Expo. SHIPPED. Real signed artifacts, v1.1.0 through 1.12.0+"),
            ("native/", "Kotlin buildable, Swift NEVER COMPILED. A candid prototype."),
            ("hardware/", "Real KiCad boards and Gerbers, generated from BOM + schematic"),
            ("Not manufactured", "Certification, tooling and fabrication remain unchecked"),
        ],
    },
    {
        "kind": "scorecard",
        "title": "Health assessment",
        "rows": [
            ("Incident documentation", 5, "Eleven quoted examples"),
            ("Security headers", 5, "Complete, applied at edge and globally"),
            ("Secret history hygiene", 5, "Provably clean across all branches"),
            ("Verification thinking", 4, "Scripts prove rather than assert"),
            ("Test volume", 4, "4,328 tests, one failing"),
            ("Hardware engineering", 4, "17 real board designs with Gerbers"),
            ("Code hygiene", 3, "0 ts-ignore, 1 TODO — but lint config broken"),
            ("Documentation accuracy", 2, "README describes an eighth of it"),
            ("Auth architecture", 2, "Five schemes, no central gate"),
            ("Payment reconciliation", 2, "Webhook verifies, then does nothing"),
            ("Schema management", 2, "Created at runtime, no migrations"),
            ("Firmware supply chain", 1, "No image signature verification"),
            ("Data durability", 1, "~27 modules memory-only in production"),
            ("Observability", 1, "No error tracking, no alerting, no sink"),
            ("CI actually running", 1, "0 of 27 runs have ever started"),
        ],
    },
    {
        "kind": "bullets",
        "title": "Top five gaps",
        "lead": "One is losing data right now. One takes a day.",
        "bullets": [
            ("~27 storage modules are memory-only in production. Passkeys and "
             "developer tokens are among them. The fix is a flag on a "
             "function that already supports it.  [3 weeks]", 0),
            ("Firmware OTA has no image signature verification, on devices "
             "controlling mains relays and door locks — and the manifest "
             "endpoint they poll does not exist.  [6 weeks]", 0),
            ("CI has never run, and verify:secrets is not a CI gate. "
             "Thirteen hard gates and 4,328 tests are decorative.  [1 day]", 0),
            ("The payment webhook verifies its signature and then does "
             "nothing, so a payment without a browser return is never "
             "reconciled.  [2 weeks]", 0),
            ("Database schema is created at runtime by the app booting — no "
             "migrations, no foreign keys, no access control of any kind.  "
             "[1 month]", 0),
        ],
    },
    {
        "kind": "roadmap",
        "title": "Phased roadmap",
        "phases": [
            ("Phase 1", "about 1 week", "Turn the lights on",
             "Fix whatever prevents GitHub Actions from starting · add "
             "verify:secrets to CI · fix the failing report-logo test · fix "
             "the eslint ignore patterns so lint reports 705 real problems "
             "instead of 29,687 · rewrite the README"),
            ("Phase 2", "about 1 month", "Stop losing data",
             "Audit all thirty createFileStore callers · set durable: true on "
             "every module whose data matters · generalise the four-process "
             "durability harness from one module to all of them"),
            ("Phase 3", "about 1 month", "Money",
             "Wire the payment webhook to actually mark orders paid or "
             "failed · add a reconciliation job against the gateway · give "
             "money a branded integer type"),
            ("Phase 4-5", "1-2 quarters", "Supply chain, then foundations",
             "Ed25519 firmware signing verified on-device before flashing · "
             "build the OTA manifest endpoint · rollback and key-rotation "
             "policy · real migrations · split store_kv into real tables · "
             "error tracking and alerting · one shared capability schema"),
        ],
    },
    {
        "kind": "bullets",
        "title": "What must NOT change",
        "lead": "Decisions that look like overhead and are not.",
        "bullets": [
            ("The incident-comment convention. Eleven examples are quoted in "
             "these documents. It made this audit possible.", 0),
            ("Verification that proves rather than asserts — four real "
             "processes, computed CSS in a real browser, a generated deck "
             "opened and checked.", 0),
            ("Separate session secrets for staff and customers, plus the "
             "token version counter added after a departing employee's copied "
             "token stayed valid forever.", 0),
            ("Passkey scopes that \u201cmust never be interchangeable\u201d, "
             "and server-side TOTP QR rendering so the secret never reaches a "
             "third-party service.", 0),
            ("Payment capture re-fetched from the gateway, and image "
             "remotePatterns scoped to two hosts: \u201c/** would turn the "
             "Next image optimizer into an open proxy.\u201d", 0),
            ("The drone companion-computer design: \u201cWHY THE CLOUD IS "
             "NEVER IN THE CONTROL LOOP... There is deliberately no 'nudge "
             "forward while I hold this button'.\u201d", 0),
        ],
    },
    {
        "kind": "closing",
        "title": "In one sentence",
        "lead": "Real engineering judgement, almost none of it automated.",
        "body": "Four hundred and twenty-eight thousand lines spanning a "
                "storefront, an admin console, an IoT platform, firmware for "
                "seventeen retail products, real PCB designs and a shipped "
                "mobile app. Modules that open by naming the production "
                "incident they exist to prevent. Verification scripts that "
                "spawn four real operating-system processes rather than "
                "mocking a cold start, because \u201c'the command succeeded' "
                "is not evidence.\u201d A complete security-header set. A git "
                "history provably free of any committed secret.\n\n"
                "And: a CI pipeline of thirteen hard gates that has never "
                "executed once in twenty-seven attempts; roughly twenty-seven "
                "storage modules — including passkeys — that lose everything "
                "on a serverless cold start; firmware updates with no "
                "signature verification on devices that switch mains relays "
                "and door locks; a payment webhook that checks its signature "
                "and then does nothing; and a README describing about an "
                "eighth of what is here.\n\n"
                "This repository knows what good looks like. It has written "
                "the reasons down, in the files themselves. What it lacks is "
                "anything that runs them without being asked.",
    },
]


def build_pptx(path: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW = prs.slide_width
    SH = prs.slide_height
    blank = prs.slide_layouts[6]

    def rgb(h):
        return RGBColor.from_string(h)

    def textbox(slide, x, y, w, h, text, size=16, bold=False, color=INK,
                align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP,
                spacing=1.0):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        lines = text.split("\n")
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.line_spacing = spacing
            r = p.add_run()
            r.text = strip_glyphs(ln)
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.name = font
            r.font.color.rgb = rgb(color)
        return tb

    def rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
        sh = slide.shapes.add_shape(shape, x, y, w, h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = rgb(fill)
        if line:
            sh.line.color.rgb = rgb(line)
            sh.line.width = Pt(1.2)
        else:
            sh.line.fill.background()
        sh.shadow.inherit = False
        sh.text_frame.word_wrap = True
        return sh

    def label(shape, text, size=12, bold=True, color="FFFFFF",
              align=PP_ALIGN.CENTER):
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, ln in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = strip_glyphs(ln)
            r.font.size = Pt(size)
            r.font.bold = bold and i == 0
            r.font.name = "Calibri"
            r.font.color.rgb = rgb(color)

    def chrome(slide, title, index):
        rect(slide, 0, 0, SW, Inches(0.92), BAND, shape=MSO_SHAPE.RECTANGLE)
        textbox(slide, Inches(0.55), Inches(0.16), SW - Inches(2.4), Inches(0.6),
                title, size=25, bold=True, color="FFFFFF")
        rect(slide, 0, Inches(0.92), SW, Emu(38100), ACCENT,
             shape=MSO_SHAPE.RECTANGLE)
        textbox(slide, SW - Inches(1.5), Inches(0.26), Inches(1.0), Inches(0.4),
                f"{index:02d}", size=15, bold=True, color="64748B",
                align=PP_ALIGN.RIGHT)
        textbox(slide, Inches(0.55), SH - Inches(0.46),
                Inches(8.0), Inches(0.3),
                f"{TITLE} · {SUBTITLE} · {BUILT}", size=9, color=MUTED)

    for i, spec in enumerate(SLIDES):
        slide = prs.slides.add_slide(blank)
        kind = spec["kind"]

        # ── title slide ──
        if kind == "title":
            rect(slide, 0, 0, SW, SH, BAND, shape=MSO_SHAPE.RECTANGLE)
            rect(slide, 0, SH - Inches(0.13), SW, Inches(0.13), ACCENT,
                 shape=MSO_SHAPE.RECTANGLE)
            textbox(slide, Inches(1.1), Inches(2.35), SW - Inches(2.2), Inches(1.3),
                    spec["title"], size=54, bold=True, color="FFFFFF",
                    align=PP_ALIGN.CENTER)
            rl = rect(slide, (SW - Inches(3.0)) // 2, Inches(3.72), Inches(3.0),
                      Emu(38100), ACCENT, shape=MSO_SHAPE.RECTANGLE)
            textbox(slide, Inches(1.1), Inches(4.05), SW - Inches(2.2), Inches(1.2),
                    spec["subtitle"], size=19, color="C7D2DE",
                    align=PP_ALIGN.CENTER, spacing=1.35)
            textbox(slide, Inches(1.1), Inches(5.72), SW - Inches(2.2), Inches(0.9),
                    "151,500 lines  ·  112 API routes  ·  "
                    "one 956 MB mail server", size=13, color="7F8EA3",
                    align=PP_ALIGN.CENTER)
            continue

        chrome(slide, spec["title"], i)
        top = Inches(1.28)

        # ── bulleted slide ──
        if kind == "bullets":
            y = top
            if spec.get("lead"):
                bar = rect(slide, Inches(0.55), y, SW - Inches(1.1), Inches(0.62),
                           "EEF3FB", ACCENT)
                label(bar, spec["lead"], size=16, color=ACCENT_DK,
                      align=PP_ALIGN.LEFT)
                bar.text_frame.margin_left = Inches(0.22)
                y += Inches(0.88)
            for text, depth in spec["bullets"]:
                h = Inches(0.52) if len(text) < 96 else Inches(0.78)
                x = Inches(0.75) + Inches(0.5) * depth
                w = SW - Inches(1.5) - Inches(0.5) * depth
                dot = rect(slide, x, y + Inches(0.13), Inches(0.12), Inches(0.12),
                           ACCENT if depth == 0 else "94A3B8", shape=MSO_SHAPE.OVAL)
                textbox(slide, x + Inches(0.3), y, w - Inches(0.3), h, text,
                        size=15 if depth == 0 else 13.5,
                        color=INK if depth == 0 else MUTED, spacing=1.12)
                y += h + Inches(0.06)

        # ── fact table ──
        elif kind == "facts":
            y = top + Inches(0.08)
            rowh = Inches(0.62)
            for n, (k, v) in enumerate(spec["facts"]):
                bg = "F1F5F9" if n % 2 == 0 else "FFFFFF"
                rect(slide, Inches(0.55), y, SW - Inches(1.1), rowh, bg,
                     shape=MSO_SHAPE.RECTANGLE)
                rect(slide, Inches(0.55), y, Emu(38100), rowh, ACCENT,
                     shape=MSO_SHAPE.RECTANGLE)
                textbox(slide, Inches(0.78), y + Inches(0.13), Inches(3.5),
                        Inches(0.4), k, size=14, bold=True, color=ACCENT_DK)
                textbox(slide, Inches(4.45), y + Inches(0.13),
                        SW - Inches(5.1), Inches(0.4), v, size=14, color=INK)
                y += rowh + Inches(0.07)

        # ── topology diagram ──
        elif kind == "topology":
            # Data-driven. spec keys:
            #   top       [(label, sub, tone)]  tone: neutral | primary | warn
            #   links     [str] connectors between adjacent top boxes
            #   hub       index of the box the children hang beneath
            #   children  [(label, sub)] up to 5
            #   childnote caption under the children row
            TONES = {
                "neutral": ("F8FAFC", "64748B", INK),
                "primary": ("DBEAFE", ACCENT, ACCENT_DK),
                "warn": ("FEF3C7", WARN, "7C2D12"),
                "good": ("ECFDF5", GOOD, "14532D"),
                "violet": ("F5F3FF", "6D28D9", "4C1D95"),
            }
            top = spec.get("top") or [("Clients", "browser", "neutral"),
                                      (TITLE, "this application", "primary")]
            links = spec.get("links") or []
            hub = spec.get("hub", len(top) - 1)
            children = spec.get("children") or []

            n_top = len(top)
            gap = Inches(1.15) if n_top > 2 else Inches(1.6)
            bw = min(Inches(2.9), (SW - Inches(1.1) - gap * (n_top - 1)) // n_top)
            bh = Inches(1.15)
            span = bw * n_top + gap * (n_top - 1)
            x0 = (SW - span) // 2
            cy = Inches(1.7)

            xs = [x0 + i * (bw + gap) for i in range(n_top)]
            for i, (lab, sub, tone) in enumerate(top):
                fill, stroke, ink = TONES.get(tone, TONES["neutral"])
                bx = rect(slide, xs[i], cy, bw, bh, fill, stroke)
                label(bx, f"{lab}\n{sub}" if sub else lab, size=14, color=ink)

            for i, txt in enumerate(links[: n_top - 1]):
                xa = xs[i] + bw
                rect(slide, xa + Inches(0.1), cy + bh // 2 - Emu(19050),
                     gap - Inches(0.2), Emu(38100), "94A3B8",
                     shape=MSO_SHAPE.RECTANGLE)
                textbox(slide, xa, cy + bh // 2 - Inches(0.44), gap, Inches(0.32),
                        txt, size=10, color=MUTED, align=PP_ALIGN.CENTER)

            if children:
                row_y = cy + bh + Inches(1.05)
                row_h = Inches(0.95)
                bus_y = cy + bh + Inches(0.52)
                sgap = Inches(0.22)
                margin = Inches(0.55)
                sw = (SW - margin * 2 - sgap * (len(children) - 1)) // len(children)
                centres = []
                for n, (lab, sub) in enumerate(children):
                    x = margin + n * (sw + sgap)
                    bx = rect(slide, x, row_y, sw, row_h, "ECFDF5", GOOD)
                    label(bx, f"{lab}\n{sub}" if sub else lab, size=12,
                          color="14532D")
                    centres.append(x + sw // 2)

                hub_cx = xs[max(0, min(hub, n_top - 1))] + bw // 2
                rect(slide, hub_cx - Emu(19050), cy + bh, Emu(38100),
                     bus_y - (cy + bh), GOOD, shape=MSO_SHAPE.RECTANGLE)
                rect(slide, centres[0], bus_y, centres[-1] - centres[0],
                     Emu(38100), GOOD, shape=MSO_SHAPE.RECTANGLE)
                for cx in centres:
                    rect(slide, cx - Emu(19050), bus_y, Emu(38100),
                         row_y - bus_y, GOOD, shape=MSO_SHAPE.RECTANGLE)

                cnote = spec.get("childnote")
                if cnote:
                    textbox(slide, margin, row_y + row_h + Inches(0.16),
                            SW - margin * 2, Inches(0.34), cnote, size=11,
                            color="14532D", align=PP_ALIGN.CENTER)

            if spec.get("note"):
                textbox(slide, Inches(0.6), SH - Inches(1.5),
                        SW - Inches(1.2), Inches(0.5), spec["note"],
                        size=13, bold=True, color=ACCENT_DK,
                        align=PP_ALIGN.CENTER)

        # ── layer stack ──
        elif kind == "layers":
            rows = [
                ("PRESENTATION", "React 19 server + client components · Tailwind v4",
                 "EFF6FF", ACCENT),
                ("APPLICATION", "App Router pages · server actions · route handlers",
                 "F5F3FF", "6D28D9"),
                ("DOMAIN — pure, no I/O",
                 "filter · format · resume · registration-validate · portal-groups",
                 "ECFDF5", GOOD),
                ("INTEGRATION", "ats.ts · candidate.ts · insights-reporter.ts",
                 "FFF7ED", WARN),
                ("PLATFORM", "Node.js on Vercel · CDN · 60-second ISR cache",
                 "F1F5F9", "64748B"),
            ]
            y = top + Inches(0.1)
            h = Inches(0.86)
            for lab, sub, fill, stroke in rows:
                bx = rect(slide, Inches(0.75), y, SW - Inches(1.5), h, fill, stroke)
                tf = bx.text_frame
                tf.margin_left = Inches(0.28)
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                r = p.add_run()
                r.text = lab
                r.font.size = Pt(15)
                r.font.bold = True
                r.font.name = "Calibri"
                r.font.color.rgb = rgb(stroke)
                p2 = tf.add_paragraph()
                p2.alignment = PP_ALIGN.LEFT
                r2 = p2.add_run()
                r2.text = strip_glyphs(sub)
                r2.font.size = Pt(11.5)
                r2.font.name = "Calibri"
                r2.font.color.rgb = rgb(MUTED)
                y += h + Inches(0.14)
            textbox(slide, Inches(0.75), SH - Inches(0.95), SW - Inches(1.5),
                    Inches(0.4), spec["note"], size=12, color=ACCENT_DK,
                    align=PP_ALIGN.CENTER)

        # ── numbered flow ──
        elif kind == "flow":
            steps = spec["steps"]
            cols = 4
            rows_n = 2
            cw = (SW - Inches(1.5) - Inches(0.28) * (cols - 1)) / cols
            ch = Inches(1.95)
            for n, (num, head, body) in enumerate(steps):
                r_i, c_i = divmod(n, cols)
                x = Inches(0.75) + c_i * (cw + Inches(0.28))
                y = top + Inches(0.25) + r_i * (ch + Inches(0.42))
                card = rect(slide, x, y, cw, ch, "FFFFFF", "CBD5E1")
                badge = rect(slide, x + Inches(0.16), y - Inches(0.2),
                             Inches(0.46), Inches(0.46), ACCENT,
                             shape=MSO_SHAPE.OVAL)
                label(badge, num, size=13, color="FFFFFF")
                textbox(slide, x + Inches(0.16), y + Inches(0.42),
                        cw - Inches(0.32), Inches(0.36), head, size=14.5,
                        bold=True, color=ACCENT_DK)
                textbox(slide, x + Inches(0.16), y + Inches(0.84),
                        cw - Inches(0.32), ch - Inches(1.0), body, size=11.5,
                        color=MUTED, spacing=1.1)

        # ── scorecard ──
        elif kind == "scorecard":
            rows_ = spec["rows"]
            # fit however many rows there are into the space below the title bar
            avail = SH - top - Inches(0.55)
            rowh = min(Inches(0.60), int(avail / max(len(rows_), 1)) - Inches(0.04))
            dot = int(min(Inches(0.28), rowh - Inches(0.22)))
            pad = int((rowh - dot) / 2)
            y = top + Inches(0.05)
            for n, (dim, score, note) in enumerate(rows_):
                bg = "F8FAFC" if n % 2 == 0 else "FFFFFF"
                rect(slide, Inches(0.55), y, SW - Inches(1.1), rowh, bg,
                     shape=MSO_SHAPE.RECTANGLE)
                textbox(slide, Inches(0.78), y + pad - Inches(0.06), Inches(3.1),
                        Inches(0.4), dim, size=14, bold=True, color=INK)
                col = GOOD if score >= 4 else (WARN if score == 3 else BAD)
                for s in range(5):
                    filled = s < score
                    rect(slide, Inches(4.0) + s * Inches(0.36),
                         y + pad, dot, dot,
                         col if filled else "E2E8F0", shape=MSO_SHAPE.OVAL)
                textbox(slide, Inches(6.15), y + pad - Inches(0.05),
                        SW - Inches(6.8), Inches(0.4), note, size=12.5,
                        color=MUTED)
                y += rowh + Inches(0.04)

        # ── roadmap ──
        elif kind == "roadmap":
            phases = spec["phases"]
            cw = (SW - Inches(1.5) - Inches(0.26) * 3) / 4
            colors_ = [(BAD, "FEF2F2"), (WARN, "FFFBEB"), (ACCENT, "EFF6FF"),
                       (GOOD, "F0FDF4")]
            for n, (name, dur, head, body) in enumerate(phases):
                x = Inches(0.75) + n * (cw + Inches(0.26))
                stroke, fill = colors_[n]
                card = rect(slide, x, top + Inches(0.35), cw, Inches(3.9),
                            fill, stroke)
                cap = rect(slide, x, top + Inches(0.35), cw, Inches(0.62),
                           stroke, shape=MSO_SHAPE.RECTANGLE)
                label(cap, f"{name}   {dur}", size=13.5, color="FFFFFF")
                textbox(slide, x + Inches(0.2), top + Inches(1.14),
                        cw - Inches(0.4), Inches(0.62), head, size=16,
                        bold=True, color=stroke)
                textbox(slide, x + Inches(0.2), top + Inches(1.86),
                        cw - Inches(0.4), Inches(2.2),
                        body.replace(" · ", "\n"), size=12.5, color=INK,
                        spacing=1.35)
            textbox(slide, Inches(0.75), SH - Inches(1.05), SW - Inches(1.5),
                    Inches(0.5),
                    "Roughly ten days of work moves operational maturity from "
                    "two stars to four.",
                    size=13, bold=True, color=ACCENT_DK, align=PP_ALIGN.CENTER)

        # ── closing ──
        elif kind == "closing":
            bar = rect(slide, Inches(0.75), top + Inches(0.2),
                       SW - Inches(1.5), Inches(0.95), "EEF3FB", ACCENT)
            label(bar, spec["lead"], size=21, color=ACCENT_DK)
            textbox(slide, Inches(1.1), top + Inches(1.55), SW - Inches(2.2),
                    Inches(3.4), spec["body"], size=15.5, color=INK,
                    spacing=1.4)

    prs.save(path)


# ─────────────────────────────────────────────────────────────── main ──

def main() -> int:
    print("=" * 66)
    print(f"  {TITLE} — documentation build")
    print("=" * 66)

    missing = [f for f, _ in SOURCES if not os.path.exists(os.path.join(HERE, f))]
    if missing:
        print("  ERROR: missing source documents: " + ", ".join(missing))
        return 1

    total = 0
    for fname, _ in SOURCES:
        size = os.path.getsize(os.path.join(HERE, fname))
        total += size
        print(f"  source   {fname:<38} {size / 1024:8.1f} KB")
    print(f"  {'':<49}{'-' * 11}")
    print(f"  {'total source':<49}{total / 1024:8.1f} KB\n")

    outputs = []

    md_path = os.path.join(HERE, "Architecture_Guide.md")
    with open(md_path, "w", encoding="utf-8") as fh:
        fh.write(build_master_md())
    outputs.append(md_path)
    print(f"  [1/4] Markdown  -> Architecture_Guide.md")

    docx_path = os.path.join(HERE, "Architecture_Guide.docx")
    build_docx(docx_path)
    outputs.append(docx_path)
    print(f"  [2/4] Word      -> Architecture_Guide.docx")

    pdf_path = os.path.join(HERE, "Architecture_Guide.pdf")
    build_pdf(pdf_path)
    outputs.append(pdf_path)
    print(f"  [3/4] PDF       -> Architecture_Guide.pdf")

    pptx_path = os.path.join(HERE, "Architecture_Overview.pptx")
    build_pptx(pptx_path)
    outputs.append(pptx_path)
    print(f"  [4/4] PowerPoint-> Architecture_Overview.pptx")

    print("\n" + "-" * 66)
    ok = True
    for p in outputs:
        if os.path.exists(p) and os.path.getsize(p) > 2048:
            print(f"  OK   {os.path.basename(p):<32} {os.path.getsize(p) / 1024:8.1f} KB")
        else:
            print(f"  FAIL {os.path.basename(p)}")
            ok = False
    print("-" * 66)
    print("  Build complete." if ok else "  Build FAILED.")
    return 0 if ok else 1


if __name__ == "__main__":
    sys.exit(main())
